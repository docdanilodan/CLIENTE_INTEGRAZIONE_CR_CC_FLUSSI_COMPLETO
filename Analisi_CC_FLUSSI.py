"""
FinancePlus.Tech - Analisi CC/FLUSSI
====================================
Applicazione Streamlit in un unico file per:
- caricare documenti di qualsiasi estensione e leggere automaticamente i principali formati bancari, PDF, Word, Excel, immagini e archivi;
- estrarre/integrare l'anagrafica aziendale;
- evitare duplicazioni tramite Partita IVA/Codice Fiscale quando presenti e tramite denominazione negli altri casi;
- calcolare KPI di liquidita, flussi, concentrazione, debito e red flags;
- generare un report PDF professionale di massimo quattro sezioni/pagine;
- archiviare, visualizzare, scaricare ed eliminare i report.

Avvio:
    pip install -r requirements_analisi_cc.txt
    streamlit run Analisi_CC_FLUSSI.py

Nota metodologica:
Lo score e la PD sono indicatori interni gestionali, non costituiscono rating
regolamentare, decisione automatizzata di credito, controllo AML o certificazione.
"""

from __future__ import annotations

import base64
import csv
import hashlib
import io
import json
import math
import os
import re
import sqlite3
import statistics
import textwrap
import uuid
import zipfile
import xml.etree.ElementTree as ET
from collections import Counter, defaultdict
from dataclasses import dataclass, asdict
from datetime import date, datetime
from pathlib import Path
from typing import Any, Iterable, Optional

import numpy as np
import pandas as pd
import streamlit as st
import streamlit.components.v1 as components

try:
    import pdfplumber
except Exception:  # pragma: no cover
    pdfplumber = None

try:
    from docx import Document
except Exception:  # pragma: no cover
    Document = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None

try:
    from striprtf.striprtf import rtf_to_text
except Exception:  # pragma: no cover
    rtf_to_text = None

try:
    import pytesseract
except Exception:  # pragma: no cover
    pytesseract = None

try:
    from PIL import Image as PILImage
except Exception:  # pragma: no cover
    PILImage = None

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm, mm
from reportlab.platypus import (
    BaseDocTemplate,
    Frame,
    Image,
    KeepTogether,
    PageBreak,
    PageTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
)
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics


# -----------------------------------------------------------------------------
# CONFIGURAZIONE
# -----------------------------------------------------------------------------
APP_TITLE = "Analisi CC/FLUSSI"
APP_SUBTITLE = "Valutazione creditizia PMI basata su estratti conto e flussi di cassa"
APP_VERSION = "1.1.0"
BRAND = "FinancePlus.Tech"
FOOTER_URL_DEFAULT = "www.financeplus.tech"

NAVY = "#0B1F3A"
NAVY_2 = "#12355B"
COPPER = "#B87333"
GOLD = "#D8B15A"
LIGHT = "#F4F7FB"
MID = "#D9E2EF"
GREEN = "#16835A"
AMBER = "#C78316"
RED = "#B63131"
TEXT = "#17202A"
MUTED = "#617083"

SCRIPT_DIR = Path(__file__).resolve().parent
DATA_DIR = SCRIPT_DIR / "analisi_cc_flussi_data"
ASSETS_DIR = DATA_DIR / "assets"
CLIENTS_DIR = DATA_DIR / "clienti"
DB_PATH = DATA_DIR / "analisi_cc_flussi.sqlite3"
SETTINGS_PATH = DATA_DIR / "settings.json"

SUPPORTED_TYPES = [
    "pdf", "csv", "tsv", "xlsx", "xls", "xlsm", "ods", "ofx", "qif", "txt", "sta",
    "xml", "json", "docx", "doc", "odt", "rtf", "html", "htm", "pptx", "zip",
    "png", "jpg", "jpeg", "tif", "tiff", "bmp", "webp"
]

DATE_HINTS = [
    "data", "date", "bookingdate", "accountingdate", "datacontabile",
    "datamovimento", "valuta", "datavaluta", "operationdate",
]
DESCRIPTION_HINTS = [
    "descrizione", "causale", "description", "memo", "dettaglio", "operazione",
    "beneficiario", "ordinante", "nome", "narrative", "transaction",
]
AMOUNT_HINTS = ["importo", "amount", "valore", "netto", "movimento", "transactionamount"]
DEBIT_HINTS = ["dare", "debit", "addebiti", "uscite", "withdrawal", "amountdebit"]
CREDIT_HINTS = ["avere", "credit", "accrediti", "entrate", "deposit", "amountcredit"]
BALANCE_HINTS = ["saldo", "balance", "disponibile", "bookbalance", "runningbalance"]

BANK_NAMES = [
    "INTESA SANPAOLO", "UNICREDIT", "BANCO BPM", "BPER BANCA", "BNL",
    "CREDIT AGRICOLE", "CREDITO EMILIANO", "CREDEM", "MONTE DEI PASCHI",
    "MPS", "FINECO", "MEDIOBANCA PREMIER", "BANCA SELLA", "ILLIMITY",
    "BANCA IFIS", "BPER", "BCC", "CASSA CENTRALE", "REVOLUT", "WISE",
    "QONTO", "N26", "POSTE ITALIANE", "BANCOPOSTA", "ING",
]

COUNTRY_RISK_KEYWORDS = [
    "CAYMAN", "PANAMA", "BELIZE", "SEYCHELLES", "VANUATU", "BVI",
    "BRITISH VIRGIN", "MARSHALL ISLAND", "SAMOA", "MAURITIUS",
]

CATEGORY_PATTERNS: list[tuple[str, str]] = [
    ("Red flag - Gioco e scommesse", r"\b(SNAI|SISAL|BET365|POKER|CASINO|SCOMMESS\w*|GIOCO ONLINE|LOTTOMATICA)\b"),
    ("Red flag - Recupero crediti/Pignoramenti", r"\b(PIGNORAMENT\w*|RECUPERO CREDITI|DECRETO INGIUNTIVO|PRECETTO|ATTO GIUDIZIARIO|UFFICIALE GIUDIZIARIO)\b"),
    ("Insoluti e respinti", r"\b(INSOLUT\w*|RESPINT\w*|STORNAT\w*|MANCANZA FONDI|NON PAGATO|IMPAGAT\w*|RIBA INSOLUTA|SDD RESPINTO|ASSEGNO IMPAGATO)\b"),
    ("Debito finanziario", r"\b(MUTUO|RATA|FINANZIAMENT\w*|LEASING|PRESTIT\w*|NOLEGGIO FINANZIARIO|CREDITO FONDIARIO|COMPASS|FINDOMESTIC|AGOS|SANTANDER CONSUMER)\b"),
    ("Anticipo fatture/Ri.Ba.", r"\b(ANTICIPO FATTUR\w*|SBF|SALVO BUON FINE|RI\.?BA\.?|RICEVUT[AE] BANCARI[AE]|SMOBILIZZ\w*|FACTORING|ANTICIPO CREDITI)\b"),
    ("Tasse e contributi", r"\b(F24|AGENZIA ENTRATE|ADER|RISCOSSIONE|IMPOST\w*|TRIBUT\w*|IVA|INPS|INAIL|TARI|IMU|IRAP|IRES|RITENUT\w*)\b"),
    ("Personale", r"\b(STIPEND\w*|EMOLUMENT\w*|SALAR\w*|RETRIBUZ\w*|PAGH\w*|DIPENDENT\w*|BONUS DIPENDENTI|TREDICESIMA|QUATTORDICESIMA)\b"),
    ("Utenze e servizi", r"\b(ENEL|ENI PLENITUDE|EDISON|A2A|HERA|ACEA|TELECOM|TIM |VODAFONE|WINDTRE|FASTWEB|ILIAD|AWS|AMAZON WEB SERVICES|MICROSOFT|GOOGLE CLOUD|UTENZ\w*|ENERGIA|GAS|ACQUA)\b"),
    ("Fornitori commerciali", r"\b(FATTUR\w*|FORNIT\w*|ACQUISTO MERCE|SALDO FATT\w*|S\.R\.L\.|SRL|S\.P\.A\.|SPA|S\.N\.C\.|SNC|SAS|COOPERATIVA)\b"),
    ("Prelievi contante", r"\b(PRELIEVO|ATM|BANCOMAT CASH|CASH WITHDRAWAL)\b"),
    ("Spese bancarie", r"\b(COMMISSION\w*|CANONE CONTO|SPESE CONTO|SPESE LIQUIDAZIONE|INTERESSI PASSIVI|BOLLO CONTO|VALUTA ADDEBITO SPESE)\b"),
    ("Affitti e immobili", r"\b(AFFITTO|CANONE LOCAZIONE|LOCAZIONE|CONDOMINIO|FITTO)\b"),
    ("Assicurazioni", r"\b(ASSICURAZ\w*|POLIZZ\w*|GENERALI|UNIPOL|ALLIANZ|AXA|ZURICH)\b"),
    ("Carte e POS", r"\b(POS|CARTA DI CREDITO|NEXI|AMERICAN EXPRESS|MASTERCARD|VISA|PAGOBANCOMAT)\b"),
    ("Giroconti/Parti correlate", r"\b(GIROCONTO|SOCIO|AMMINISTRATORE|PRELEVAMENTO SOCI|FINANZIAMENTO SOCI|CONTO PERSONALE|DIVIDENDO)\b"),
]


# -----------------------------------------------------------------------------
# DATA CLASS
# -----------------------------------------------------------------------------
@dataclass
class ClientInfo:
    legal_name: str = ""
    vat_number: str = ""
    tax_code: str = ""
    address: str = ""
    city: str = ""
    province: str = ""
    postal_code: str = ""
    ateco: str = ""
    bank_name: str = ""
    iban: str = ""


# -----------------------------------------------------------------------------
# UTILITA GENERALI
# -----------------------------------------------------------------------------
def ensure_directories() -> None:
    for folder in (DATA_DIR, ASSETS_DIR, CLIENTS_DIR):
        folder.mkdir(parents=True, exist_ok=True)


def now_iso() -> str:
    return datetime.now().replace(microsecond=0).isoformat()


def now_display() -> str:
    return datetime.now().strftime("%d-%m-%Y %H:%M")


def safe_filename(value: str, fallback: str = "SENZA_NOME") -> str:
    value = (value or fallback).strip().upper()
    value = re.sub(r"[^A-Z0-9À-ÖØ-Ý]+", "_", value)
    value = re.sub(r"_+", "_", value).strip("_")
    return value[:100] or fallback


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def combined_hash(files: list[dict[str, Any]]) -> str:
    h = hashlib.sha256()
    for item in sorted(files, key=lambda x: x["name"]):
        h.update(item["name"].encode("utf-8", errors="ignore"))
        h.update(item["bytes"])
    return h.hexdigest()


def euro(value: Any, empty: str = "-") -> str:
    try:
        if value is None or (isinstance(value, float) and math.isnan(value)):
            return empty
        return f"€ {float(value):,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    except Exception:
        return empty


def pct(value: Any, digits: int = 1, empty: str = "-") -> str:
    try:
        if value is None or (isinstance(value, float) and math.isnan(value)):
            return empty
        return f"{float(value) * 100:.{digits}f}%".replace(".", ",")
    except Exception:
        return empty


def num_it(value: Any, digits: int = 2, empty: str = "-") -> str:
    try:
        if value is None or (isinstance(value, float) and math.isnan(value)):
            return empty
        return f"{float(value):,.{digits}f}".replace(",", "X").replace(".", ",").replace("X", ".")
    except Exception:
        return empty


def clean_text(value: Any) -> str:
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return ""
    text = str(value).replace("\x00", " ")
    return re.sub(r"\s+", " ", text).strip()


def normalize_col(value: Any) -> str:
    text = clean_text(value).lower()
    text = text.replace("à", "a").replace("è", "e").replace("é", "e").replace("ì", "i").replace("ò", "o").replace("ù", "u")
    return re.sub(r"[^a-z0-9]+", "", text)


def parse_amount(value: Any) -> float:
    """Converte importi italiani/internazionali in float."""
    if value is None:
        return float("nan")
    if isinstance(value, (int, float, np.number)):
        return float(value)
    s = str(value).strip()
    if not s or s.lower() in {"nan", "none", "-", "--"}:
        return float("nan")
    neg = False
    if s.startswith("(") and s.endswith(")"):
        neg = True
        s = s[1:-1]
    s = s.replace("€", "").replace("EUR", "").replace("eur", "")
    s = re.sub(r"[^0-9,\.\-+]", "", s)
    if not s:
        return float("nan")

    # Determina il separatore decimale dall'ultima occorrenza.
    if "," in s and "." in s:
        if s.rfind(",") > s.rfind("."):
            s = s.replace(".", "").replace(",", ".")
        else:
            s = s.replace(",", "")
    elif "," in s:
        parts = s.split(",")
        if len(parts[-1]) in (1, 2, 3):
            s = "".join(parts[:-1]) + "." + parts[-1]
        else:
            s = s.replace(",", "")
    elif s.count(".") > 1:
        parts = s.split(".")
        if len(parts[-1]) in (1, 2):
            s = "".join(parts[:-1]) + "." + parts[-1]
        else:
            s = s.replace(".", "")

    try:
        result = float(s)
        return -abs(result) if neg else result
    except ValueError:
        return float("nan")


def parse_date_series(series: pd.Series) -> pd.Series:
    parsed = pd.to_datetime(series, errors="coerce", dayfirst=True)
    if parsed.notna().mean() < 0.35:
        parsed = pd.to_datetime(series, errors="coerce", dayfirst=False)
    return parsed


def coalesce(*values: Any) -> Any:
    for value in values:
        if value not in (None, "") and not (isinstance(value, float) and math.isnan(value)):
            return value
    return ""


def load_settings() -> dict[str, Any]:
    defaults = {
        "footer_url": FOOTER_URL_DEFAULT,
        "consultant_name": "FinancePlus.Tech",
        "default_credit_limit": 0.0,
        "default_requested_line": 0.0,
    }
    if SETTINGS_PATH.exists():
        try:
            data = json.loads(SETTINGS_PATH.read_text(encoding="utf-8"))
            defaults.update(data)
        except Exception:
            pass
    return defaults


def save_settings(settings: dict[str, Any]) -> None:
    SETTINGS_PATH.write_text(json.dumps(settings, ensure_ascii=False, indent=2), encoding="utf-8")


def find_logo_path() -> Optional[Path]:
    candidates = [
        ASSETS_DIR / "logo_financeplus.png",
        SCRIPT_DIR / "LOGO1.png",
        SCRIPT_DIR / "LOGO_FINANCE_2.PNG",
        SCRIPT_DIR / "logo_financeplus.png",
        SCRIPT_DIR / "logo_banner.png",
        SCRIPT_DIR / "logo.png",
    ]
    for path in candidates:
        if path.exists() and path.is_file():
            return path
    return None


def register_fonts() -> tuple[str, str]:
    """Registra DejaVu se disponibile, altrimenti usa Helvetica."""
    regular = "Helvetica"
    bold = "Helvetica-Bold"
    candidates = [
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/Library/Fonts/Arial Unicode.ttf"),
        Path("C:/Windows/Fonts/arial.ttf"),
    ]
    bold_candidates = [
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        Path("/Library/Fonts/Arial Bold.ttf"),
        Path("C:/Windows/Fonts/arialbd.ttf"),
    ]
    try:
        reg_path = next((p for p in candidates if p.exists()), None)
        bold_path = next((p for p in bold_candidates if p.exists()), None)
        if reg_path:
            pdfmetrics.registerFont(TTFont("FPRegular", str(reg_path)))
            regular = "FPRegular"
        if bold_path:
            pdfmetrics.registerFont(TTFont("FPBold", str(bold_path)))
            bold = "FPBold"
    except Exception:
        pass
    return regular, bold


PDF_FONT, PDF_FONT_BOLD = register_fonts()


# -----------------------------------------------------------------------------
# DATABASE
# -----------------------------------------------------------------------------
def db_connect() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    return conn


def init_db() -> None:
    ensure_directories()
    with db_connect() as conn:
        conn.executescript(
            """
            CREATE TABLE IF NOT EXISTS clients (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                legal_name TEXT NOT NULL,
                vat_number TEXT,
                tax_code TEXT,
                address TEXT,
                city TEXT,
                province TEXT,
                postal_code TEXT,
                ateco TEXT,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            );

            CREATE UNIQUE INDEX IF NOT EXISTS idx_clients_vat
            ON clients(vat_number) WHERE vat_number IS NOT NULL AND vat_number <> '';

            CREATE UNIQUE INDEX IF NOT EXISTS idx_clients_tax_code
            ON clients(tax_code) WHERE tax_code IS NOT NULL AND tax_code <> '';

            CREATE TABLE IF NOT EXISTS reports (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                client_id INTEGER NOT NULL,
                report_code TEXT NOT NULL UNIQUE,
                version INTEGER NOT NULL DEFAULT 1,
                bank_name TEXT,
                iban TEXT,
                period_start TEXT,
                period_end TEXT,
                period_option TEXT,
                source_names TEXT,
                source_hash TEXT,
                source_dir TEXT,
                pdf_path TEXT NOT NULL,
                score REAL,
                rating_class TEXT,
                risk_level TEXT,
                estimated_pd REAL,
                requested_line REAL,
                suggested_line REAL,
                summary TEXT,
                metrics_json TEXT,
                created_at TEXT NOT NULL,
                FOREIGN KEY(client_id) REFERENCES clients(id) ON DELETE CASCADE
            );

            CREATE INDEX IF NOT EXISTS idx_reports_client ON reports(client_id);
            CREATE INDEX IF NOT EXISTS idx_reports_date ON reports(created_at);
            CREATE INDEX IF NOT EXISTS idx_reports_bank_period
            ON reports(client_id, bank_name, period_start, period_end);
            """
        )


def find_client(vat_number: str, tax_code: str, legal_name: str = "") -> Optional[sqlite3.Row]:
    with db_connect() as conn:
        if vat_number:
            row = conn.execute("SELECT * FROM clients WHERE vat_number = ?", (vat_number,)).fetchone()
            if row:
                return row
        if tax_code:
            row = conn.execute("SELECT * FROM clients WHERE tax_code = ?", (tax_code,)).fetchone()
            if row:
                return row
        if legal_name:
            normalized = safe_filename(legal_name)
            rows = conn.execute("SELECT * FROM clients").fetchall()
            for row in rows:
                if safe_filename(row["legal_name"]) == normalized:
                    return row
    return None


def upsert_client(client: ClientInfo) -> int:
    existing = find_client(client.vat_number, client.tax_code, client.legal_name)
    with db_connect() as conn:
        if existing:
            conn.execute(
                """
                UPDATE clients SET legal_name=?, vat_number=?, tax_code=?, address=?, city=?,
                    province=?, postal_code=?, ateco=?, updated_at=? WHERE id=?
                """,
                (
                    client.legal_name or existing["legal_name"],
                    client.vat_number or existing["vat_number"],
                    client.tax_code or existing["tax_code"],
                    client.address or existing["address"],
                    client.city or existing["city"],
                    client.province or existing["province"],
                    client.postal_code or existing["postal_code"],
                    client.ateco or existing["ateco"],
                    now_iso(),
                    existing["id"],
                ),
            )
            return int(existing["id"])
        cur = conn.execute(
            """
            INSERT INTO clients (
                legal_name, vat_number, tax_code, address, city, province,
                postal_code, ateco, created_at, updated_at
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                client.legal_name or "Cliente non identificato",
                client.vat_number,
                client.tax_code,
                client.address,
                client.city,
                client.province,
                client.postal_code,
                client.ateco,
                now_iso(),
                now_iso(),
            ),
        )
        return int(cur.lastrowid)


def count_matching_reports(client_id: int, bank_name: str, period_start: str, period_end: str) -> int:
    with db_connect() as conn:
        row = conn.execute(
            """
            SELECT COUNT(*) AS n FROM reports
            WHERE client_id=? AND UPPER(COALESCE(bank_name,''))=UPPER(?)
              AND COALESCE(period_start,'')=? AND COALESCE(period_end,'')=?
            """,
            (client_id, bank_name or "", period_start or "", period_end or ""),
        ).fetchone()
        return int(row["n"])


def insert_report(record: dict[str, Any]) -> int:
    columns = ", ".join(record.keys())
    placeholders = ", ".join(["?"] * len(record))
    with db_connect() as conn:
        cur = conn.execute(
            f"INSERT INTO reports ({columns}) VALUES ({placeholders})",
            tuple(record.values()),
        )
        return int(cur.lastrowid)


def list_reports_df() -> pd.DataFrame:
    query = """
        SELECT r.*, c.legal_name, c.vat_number, c.tax_code, c.city, c.province
        FROM reports r
        JOIN clients c ON c.id = r.client_id
        ORDER BY r.created_at DESC
    """
    with db_connect() as conn:
        return pd.read_sql_query(query, conn)


def dashboard_stats() -> dict[str, Any]:
    with db_connect() as conn:
        clients = conn.execute("SELECT COUNT(*) n FROM clients").fetchone()["n"]
        reports = conn.execute("SELECT COUNT(*) n FROM reports").fetchone()["n"]
        avg_score = conn.execute("SELECT AVG(score) v FROM reports").fetchone()["v"]
        high_risk = conn.execute("SELECT COUNT(*) n FROM reports WHERE risk_level='ALTO'").fetchone()["n"]
    return {"clients": clients, "reports": reports, "avg_score": avg_score or 0.0, "high_risk": high_risk}


def delete_report(report_id: int) -> None:
    with db_connect() as conn:
        row = conn.execute("SELECT * FROM reports WHERE id=?", (report_id,)).fetchone()
        if not row:
            return
        paths = [row["pdf_path"], row["source_dir"]]
        conn.execute("DELETE FROM reports WHERE id=?", (report_id,))
    for raw in paths:
        if not raw:
            continue
        path = Path(raw)
        try:
            if path.is_file():
                path.unlink(missing_ok=True)
            elif path.is_dir():
                for child in path.glob("**/*"):
                    if child.is_file():
                        child.unlink(missing_ok=True)
                for child in sorted(path.glob("**/*"), reverse=True):
                    if child.is_dir():
                        child.rmdir()
                path.rmdir()
        except Exception:
            pass


# -----------------------------------------------------------------------------
# ESTRAZIONE ANAGRAFICA
# -----------------------------------------------------------------------------
def first_match(patterns: Iterable[str], text: str, flags: int = re.I) -> str:
    for pattern in patterns:
        match = re.search(pattern, text, flags)
        if match:
            return clean_text(match.group(1))
    return ""


def extract_metadata(text: str, file_names: list[str]) -> ClientInfo:
    compact = re.sub(r"[\t ]+", " ", text or "")
    upper = compact.upper()

    vat = first_match(
        [
            r"(?:PARTITA\s+IVA|P\.?\s*IVA|VAT)\s*[:\-]?\s*(?:IT)?\s*(\d{11})",
            r"\bIT\s*(\d{11})\b",
        ],
        upper,
    )
    tax_code = first_match(
        [
            r"(?:CODICE\s+FISCALE|C\.?F\.?)\s*[:\-]?\s*([A-Z0-9]{11,16})",
        ],
        upper,
    )
    legal_name = first_match(
        [
            r"(?:RAGIONE\s+SOCIALE|DENOMINAZIONE|INTESTATARIO|TITOLARE\s+CONTO|CLIENTE)\s*[:\-]?\s*([^\n\r]{3,100})",
            r"(?:CONTO\s+INTESTATO\s+A)\s*[:\-]?\s*([^\n\r]{3,100})",
        ],
        compact,
    )
    if legal_name:
        legal_name = re.split(r"\s{2,}|\bIBAN\b|\bCODICE\b|\bINDIRIZZO\b", legal_name, maxsplit=1, flags=re.I)[0].strip(" -:;")

    # Fallback dal nome file.
    if not legal_name and file_names:
        stem = Path(file_names[0]).stem
        stem = re.sub(r"(?i)estratto|conto|movimenti|e[cC]|report|banca|\d{4,}", " ", stem)
        stem = re.sub(r"[_\-]+", " ", stem)
        candidate = clean_text(stem)
        if len(candidate) >= 3:
            legal_name = candidate.upper()

    address = first_match(
        [
            r"(?:INDIRIZZO|SEDE\s+LEGALE|DOMICILIO)\s*[:\-]?\s*([^\n\r]{4,120})",
            r"\b((?:VIA|VIALE|PIAZZA|CORSO|LARGO|STRADA)\s+[^\n\r]{3,100})",
        ],
        compact,
    )
    postal_code = first_match([r"\b(\d{5})\b"], address or compact)
    city_prov = re.search(r"\b([A-ZÀ-ÖØ-Ý' ]{2,40})\s*\(([A-Z]{2})\)\b", upper)
    city = clean_text(city_prov.group(1)).title() if city_prov else ""
    province = city_prov.group(2) if city_prov else ""
    ateco = first_match([r"(?:ATECO|CODICE\s+ATTIVIT[ÀA])\s*[:\-]?\s*([0-9\.]{2,10})"], upper)
    iban = first_match([r"\b(IT\d{2}[A-Z]\d{10}[A-Z0-9]{12})\b"], re.sub(r"\s+", "", upper))

    bank_name = ""
    for bank in BANK_NAMES:
        if bank in upper:
            bank_name = bank.title()
            break

    return ClientInfo(
        legal_name=legal_name,
        vat_number=vat,
        tax_code=tax_code,
        address=address,
        city=city,
        province=province,
        postal_code=postal_code,
        ateco=ateco,
        bank_name=bank_name,
        iban=iban,
    )


# -----------------------------------------------------------------------------
# PARSER DOCUMENTI
# -----------------------------------------------------------------------------
def make_unique_columns(columns: list[Any]) -> list[str]:
    seen: Counter[str] = Counter()
    result: list[str] = []
    for idx, col in enumerate(columns):
        base = clean_text(col) or f"col_{idx + 1}"
        seen[base] += 1
        result.append(base if seen[base] == 1 else f"{base}_{seen[base]}")
    return result


def parse_csv_bytes(data: bytes) -> pd.DataFrame:
    encodings = ["utf-8-sig", "utf-8", "cp1252", "latin1"]
    last_error: Optional[Exception] = None
    for enc in encodings:
        try:
            text = data.decode(enc)
            try:
                dialect = csv.Sniffer().sniff(text[:5000], delimiters=";,\t|")
                sep = dialect.delimiter
            except Exception:
                sep = None
            return pd.read_csv(io.StringIO(text), sep=sep, engine="python", dtype=str)
        except Exception as exc:
            last_error = exc
    raise ValueError(f"CSV non leggibile: {last_error}")


def parse_excel_bytes(data: bytes, extension: str) -> pd.DataFrame:
    engine = "openpyxl" if extension in {".xlsx", ".xlsm"} else ("odf" if extension == ".ods" else "xlrd")
    sheets = pd.read_excel(io.BytesIO(data), sheet_name=None, dtype=str, engine=engine)
    frames = []
    for name, frame in sheets.items():
        if frame.empty:
            continue
        frame = frame.copy()
        frame["__sheet__"] = name
        frames.append(frame)
    if not frames:
        return pd.DataFrame()
    return pd.concat(frames, ignore_index=True, sort=False)


def parse_json_bytes(data: bytes) -> pd.DataFrame:
    obj = json.loads(data.decode("utf-8-sig"))
    if isinstance(obj, dict):
        for key in ("transactions", "movimenti", "data", "items"):
            if isinstance(obj.get(key), list):
                obj = obj[key]
                break
        else:
            obj = [obj]
    return pd.json_normalize(obj)


def _xml_local_name(tag: str) -> str:
    return tag.split("}")[-1] if "}" in tag else tag


def parse_xml_bytes(data: bytes) -> tuple[pd.DataFrame, str]:
    """Parser euristico per XML bancari, inclusi export CAMT.052/053/054."""
    root = ET.fromstring(data)
    text_parts = [clean_text(elem.text) for elem in root.iter() if clean_text(elem.text)]
    text = "\n".join(text_parts)
    rows: list[dict[str, Any]] = []

    entry_names = {"Ntry", "Entry", "Transaction", "Tx", "Movement", "Movimento"}
    for entry in root.iter():
        if _xml_local_name(entry.tag) not in entry_names:
            continue
        values: dict[str, list[str]] = defaultdict(list)
        for elem in entry.iter():
            name = _xml_local_name(elem.tag)
            value = clean_text(elem.text)
            if value:
                values[name].append(value)

        date_value = coalesce(
            *(values.get("BookgDt", []) + values.get("Dt", []) + values.get("DtTm", []) +
              values.get("ValDt", []) + values.get("BookingDate", []) + values.get("Date", []))
        )
        amount_value = coalesce(*(values.get("Amt", []) + values.get("Amount", []) + values.get("TxAmt", [])))
        amount = parse_amount(amount_value)
        indicator = " ".join(values.get("CdtDbtInd", []) + values.get("CreditDebitIndicator", [])).upper()
        if not math.isnan(amount) and ("DBIT" in indicator or "DEBIT" in indicator):
            amount = -abs(amount)
        elif not math.isnan(amount) and ("CRDT" in indicator or "CREDIT" in indicator):
            amount = abs(amount)

        description_parts = []
        for key in ("Ustrd", "AddtlNtryInf", "RmtInf", "Nm", "Name", "Memo", "Description", "EndToEndId"):
            description_parts.extend(values.get(key, []))
        description = " - ".join(dict.fromkeys(description_parts))
        if date_value and not math.isnan(amount):
            rows.append({"date": date_value, "description": description, "amount": amount})

    if not rows:
        # Fallback: tenta di interpretare i nodi figli di primo/secondo livello come record tabellari.
        candidates = []
        for parent in root.iter():
            children = list(parent)
            if len(children) >= 2 and all(len(list(child)) == 0 for child in children):
                record = {_xml_local_name(child.tag): clean_text(child.text) for child in children}
                if record:
                    candidates.append(record)
        frame = pd.DataFrame(candidates)
    else:
        frame = pd.DataFrame(rows)
    return frame, text


def parse_ofx_bytes(data: bytes) -> pd.DataFrame:
    text = data.decode("utf-8", errors="ignore")
    rows = []
    for block in re.findall(r"<STMTTRN>(.*?)(?:</STMTTRN>|(?=<STMTTRN>)|$)", text, flags=re.I | re.S):
        def tag(name: str) -> str:
            m = re.search(rf"<{name}>([^<\r\n]+)", block, flags=re.I)
            return clean_text(m.group(1)) if m else ""
        rows.append({
            "date": tag("DTPOSTED")[:8],
            "amount": tag("TRNAMT"),
            "description": " - ".join(filter(None, [tag("NAME"), tag("MEMO"), tag("CHECKNUM")])),
            "type": tag("TRNTYPE"),
        })
    return pd.DataFrame(rows)


def parse_qif_bytes(data: bytes) -> pd.DataFrame:
    text = data.decode("utf-8", errors="ignore")
    records = text.split("^")
    rows = []
    for rec in records:
        fields: dict[str, str] = {}
        for line in rec.splitlines():
            if not line:
                continue
            key, value = line[0], line[1:].strip()
            fields[key] = f"{fields.get(key, '')} {value}".strip()
        if "D" in fields or "T" in fields:
            rows.append({
                "date": fields.get("D", ""),
                "amount": fields.get("T", ""),
                "description": " - ".join(filter(None, [fields.get("P", ""), fields.get("M", "")])),
            })
    return pd.DataFrame(rows)


def parse_docx_bytes(data: bytes) -> tuple[pd.DataFrame, str]:
    if Document is None:
        raise RuntimeError("python-docx non installato")
    doc = Document(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs)
    frames = []
    for table in doc.tables:
        rows = [[clean_text(cell.text) for cell in row.cells] for row in table.rows]
        if len(rows) >= 2:
            frames.append(pd.DataFrame(rows[1:], columns=make_unique_columns(rows[0])))
    return (pd.concat(frames, ignore_index=True, sort=False) if frames else pd.DataFrame(), text)


def parse_pdf_bytes(data: bytes) -> tuple[pd.DataFrame, str]:
    if pdfplumber is None:
        raise RuntimeError("pdfplumber non installato")
    frames: list[pd.DataFrame] = []
    texts: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text(x_tolerance=2, y_tolerance=3) or ""
            texts.append(page_text)
            try:
                tables = page.extract_tables() or []
            except Exception:
                tables = []
            for table in tables:
                cleaned = [[clean_text(cell) for cell in row] for row in table if row]
                cleaned = [row for row in cleaned if any(row)]
                if len(cleaned) < 2:
                    continue
                header = make_unique_columns(cleaned[0])
                width = len(header)
                body = [row[:width] + [""] * max(0, width - len(row)) for row in cleaned[1:]]
                frames.append(pd.DataFrame(body, columns=header))
    text = "\n".join(texts)
    if frames:
        return pd.concat(frames, ignore_index=True, sort=False), text
    return parse_transactions_from_text(text), text


def parse_transactions_from_text(text: str) -> pd.DataFrame:
    """Fallback euristico per PDF/TXT con una riga per movimento."""
    rows = []
    date_pattern = re.compile(r"\b(\d{1,2}[\./-]\d{1,2}[\./-](?:\d{2}|\d{4}))\b")
    amount_pattern = re.compile(r"(?<!\d)([-+]?\s*\d{1,3}(?:\.\d{3})*(?:,\d{2})|[-+]?\s*\d+(?:[\.,]\d{2}))(?!\d)")
    for line in text.splitlines():
        line = clean_text(line)
        dm = date_pattern.search(line)
        if not dm:
            continue
        amounts = amount_pattern.findall(line)
        if not amounts:
            continue
        parsed_amounts = [parse_amount(a.replace(" ", "")) for a in amounts]
        parsed_amounts = [a for a in parsed_amounts if not math.isnan(a)]
        if not parsed_amounts:
            continue
        amount = parsed_amounts[-1]
        description = line.replace(dm.group(0), " ")
        for a in amounts:
            description = description.replace(a, " ")
        rows.append({"date": dm.group(1), "description": clean_text(description), "amount": amount})
    return pd.DataFrame(rows)


def parse_rtf_bytes(data: bytes) -> tuple[pd.DataFrame, str]:
    raw = data.decode("utf-8", errors="ignore")
    text = rtf_to_text(raw) if rtf_to_text is not None else re.sub(r"\\[a-z]+\d* ?|[{}]", " ", raw)
    return parse_transactions_from_text(text), text


def parse_html_bytes(data: bytes) -> tuple[pd.DataFrame, str]:
    text = data.decode("utf-8", errors="ignore")
    try:
        tables = pd.read_html(io.StringIO(text))
        frame = pd.concat(tables, ignore_index=True, sort=False) if tables else pd.DataFrame()
    except Exception:
        frame = pd.DataFrame()
    clean = re.sub(r"<[^>]+>", " ", text)
    clean = re.sub(r"\s+", " ", clean)
    if frame.empty:
        frame = parse_transactions_from_text(clean)
    return frame, clean


def parse_pptx_bytes(data: bytes) -> tuple[pd.DataFrame, str]:
    if Presentation is None:
        raise RuntimeError("python-pptx non installato")
    prs = Presentation(io.BytesIO(data))
    texts = []
    frames = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                rows = [[clean_text(c.text) for c in row.cells] for row in shape.table.rows]
                if len(rows) >= 2:
                    frames.append(pd.DataFrame(rows[1:], columns=make_unique_columns(rows[0])))
    text = "\n".join(texts)
    frame = pd.concat(frames, ignore_index=True, sort=False) if frames else parse_transactions_from_text(text)
    return frame, text


def parse_image_bytes(data: bytes) -> tuple[pd.DataFrame, str]:
    if PILImage is None:
        raise RuntimeError("Pillow non installato")
    image = PILImage.open(io.BytesIO(data)).convert("RGB")
    if pytesseract is None:
        return pd.DataFrame(), ""
    try:
        text = pytesseract.image_to_string(image, lang="ita+eng")
    except Exception:
        text = pytesseract.image_to_string(image)
    return parse_transactions_from_text(text), text


def parse_generic_bytes(data: bytes) -> tuple[pd.DataFrame, str]:
    # Consente comunque il caricamento e l'archiviazione di estensioni non note.
    # Se il contenuto e' testuale, prova a ricostruire i movimenti.
    for enc in ("utf-8-sig", "utf-8", "cp1252", "latin1"):
        try:
            text = data.decode(enc)
            break
        except Exception:
            text = ""
    return parse_transactions_from_text(text), text


def parse_document(name: str, data: bytes) -> tuple[pd.DataFrame, str, str]:
    ext = Path(name).suffix.lower()
    text = ""
    if ext == ".csv":
        frame = parse_csv_bytes(data)
    elif ext in {".xlsx", ".xls", ".xlsm", ".ods"}:
        frame = parse_excel_bytes(data, ext)
    elif ext == ".json":
        frame = parse_json_bytes(data)
        text = json.dumps(json.loads(data.decode("utf-8-sig")), ensure_ascii=False)[:100_000]
    elif ext == ".xml":
        frame, text = parse_xml_bytes(data)
    elif ext == ".ofx":
        frame = parse_ofx_bytes(data)
        text = data.decode("utf-8", errors="ignore")
    elif ext == ".qif":
        frame = parse_qif_bytes(data)
        text = data.decode("utf-8", errors="ignore")
    elif ext == ".pdf":
        frame, text = parse_pdf_bytes(data)
    elif ext == ".docx":
        frame, text = parse_docx_bytes(data)
    elif ext == ".rtf":
        frame, text = parse_rtf_bytes(data)
    elif ext in {".html", ".htm"}:
        frame, text = parse_html_bytes(data)
    elif ext == ".pptx":
        frame, text = parse_pptx_bytes(data)
    elif ext in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}:
        frame, text = parse_image_bytes(data)
    elif ext in {".txt", ".sta", ".tsv", ".odt", ".doc"}:
        text = data.decode("utf-8", errors="ignore")
        try:
            frame = parse_csv_bytes(data)
        except Exception:
            frame = parse_transactions_from_text(text)
    else:
        frame, text = parse_generic_bytes(data)
    return frame, text, ext


def choose_column(columns: list[str], hints: list[str]) -> Optional[str]:
    normalized = {col: normalize_col(col) for col in columns}
    for hint in hints:
        for col, norm in normalized.items():
            if norm == hint or hint in norm:
                return col
    return None


def header_row_recovery(df: pd.DataFrame) -> pd.DataFrame:
    """Prova a usare una delle prime righe come intestazione quando il PDF è irregolare."""
    if df.empty:
        return df
    current_hits = sum(
        any(h in normalize_col(c) for h in DATE_HINTS + DESCRIPTION_HINTS + AMOUNT_HINTS + DEBIT_HINTS + CREDIT_HINTS)
        for c in df.columns
    )
    if current_hits >= 2:
        return df
    for idx in range(min(8, len(df))):
        values = [clean_text(v) for v in df.iloc[idx].tolist()]
        hits = sum(any(h in normalize_col(v) for h in DATE_HINTS + DESCRIPTION_HINTS + AMOUNT_HINTS + DEBIT_HINTS + CREDIT_HINTS) for v in values)
        if hits >= 2:
            out = df.iloc[idx + 1:].copy()
            out.columns = make_unique_columns(values)
            return out.reset_index(drop=True)
    return df


def standardize_transactions(raw: pd.DataFrame, source_name: str) -> tuple[pd.DataFrame, dict[str, Any]]:
    quality: dict[str, Any] = {"source": source_name, "rows_raw": len(raw), "warnings": []}
    if raw is None or raw.empty:
        quality["warnings"].append("Nessuna tabella/movimento riconosciuto")
        return pd.DataFrame(columns=["date", "description", "amount", "balance", "source"]), quality

    df = raw.copy()
    if isinstance(df.columns, pd.MultiIndex):
        df.columns = [" ".join(map(str, col)).strip() for col in df.columns]
    df.columns = make_unique_columns(list(df.columns))
    df = header_row_recovery(df)
    cols = list(df.columns)

    date_col = choose_column(cols, DATE_HINTS)
    desc_col = choose_column(cols, DESCRIPTION_HINTS)
    amount_col = choose_column(cols, AMOUNT_HINTS)
    debit_col = choose_column(cols, DEBIT_HINTS)
    credit_col = choose_column(cols, CREDIT_HINTS)
    balance_col = choose_column(cols, BALANCE_HINTS)

    # Euristica per date se intestazioni non informative.
    if date_col is None:
        best = None
        best_rate = 0.0
        for col in cols:
            rate = parse_date_series(df[col].head(100)).notna().mean()
            if rate > best_rate:
                best, best_rate = col, rate
        if best_rate >= 0.45:
            date_col = best

    # Euristica per importi se intestazioni non informative.
    if amount_col is None and not (debit_col or credit_col):
        best = None
        best_rate = 0.0
        for col in cols:
            numeric = df[col].head(200).map(parse_amount)
            rate = numeric.notna().mean()
            if rate > best_rate and col != date_col:
                best, best_rate = col, rate
        if best_rate >= 0.55:
            amount_col = best

    if desc_col is None:
        object_candidates = [c for c in cols if c not in {date_col, amount_col, debit_col, credit_col, balance_col}]
        if object_candidates:
            desc_col = max(object_candidates, key=lambda c: df[c].astype(str).str.len().mean())

    if date_col is None:
        quality["warnings"].append("Colonna data non individuata")
    if amount_col is None and not (debit_col or credit_col):
        quality["warnings"].append("Colonna importo/addebito/accredito non individuata")
    if date_col is None or (amount_col is None and not (debit_col or credit_col)):
        return pd.DataFrame(columns=["date", "description", "amount", "balance", "source"]), quality

    out = pd.DataFrame()
    out["date"] = parse_date_series(df[date_col])
    out["description"] = df[desc_col].map(clean_text) if desc_col else "Movimento bancario"

    if amount_col:
        out["amount"] = df[amount_col].map(parse_amount)
    else:
        debit = df[debit_col].map(parse_amount) if debit_col else pd.Series(0.0, index=df.index)
        credit = df[credit_col].map(parse_amount) if credit_col else pd.Series(0.0, index=df.index)
        debit = debit.fillna(0.0).abs()
        credit = credit.fillna(0.0).abs()
        out["amount"] = credit - debit

    out["balance"] = df[balance_col].map(parse_amount) if balance_col else np.nan
    out["source"] = source_name
    out = out.dropna(subset=["date", "amount"])
    out = out[out["amount"].abs() > 0.00001]
    out = out.sort_values("date").reset_index(drop=True)

    # Corregge il segno in presenza di colonne testuali Dare/Avere nella descrizione.
    if amount_col and (out["amount"] >= 0).all():
        desc_upper = out["description"].str.upper()
        likely_out = desc_upper.str.contains(r"ADDEBIT|PAGAMENTO|BONIFICO A|PRELIEVO|SDD|F24|RATA|COMMISSION", regex=True)
        likely_in = desc_upper.str.contains(r"ACCREDITO|BONIFICO DA|INCASSO|VERSAMENTO|STIPENDIO RICEVUTO", regex=True)
        if likely_out.mean() > 0.10:
            out.loc[likely_out & ~likely_in, "amount"] *= -1
            quality["warnings"].append("Segno Dare/Avere ricostruito dalla causale")

    quality.update({
        "rows_valid": len(out),
        "date_col": date_col,
        "description_col": desc_col,
        "amount_col": amount_col,
        "debit_col": debit_col,
        "credit_col": credit_col,
        "balance_col": balance_col,
    })
    if len(out) < 5:
        quality["warnings"].append("Numero di movimenti riconosciuti molto basso")
    return out, quality


def parse_uploaded_files(uploaded_files: list[Any]) -> tuple[pd.DataFrame, str, list[dict[str, Any]], list[dict[str, Any]]]:
    standardized_frames: list[pd.DataFrame] = []
    all_text: list[str] = []
    qualities: list[dict[str, Any]] = []
    source_files: list[dict[str, Any]] = []

    for uploaded in uploaded_files:
        data = uploaded.getvalue()
        source_files.append({"name": uploaded.name, "bytes": data, "hash": sha256_bytes(data)})
        if Path(uploaded.name).suffix.lower() == ".zip":
            try:
                with zipfile.ZipFile(io.BytesIO(data)) as archive:
                    members = [m for m in archive.infolist() if not m.is_dir()]
                    if len(members) > 100:
                        raise ValueError("Archivio ZIP con oltre 100 file: suddividere il caricamento")
                    total_uncompressed = sum(m.file_size for m in members)
                    if total_uncompressed > 80 * 1024 * 1024:
                        raise ValueError("Archivio ZIP oltre 80 MB non compresso")
                    parsed_members = 0
                    for member in members:
                        member_ext = Path(member.filename).suffix.lower().lstrip(".")
                        if member.endswith("/") or member_ext == "zip":
                            continue
                        member_data = archive.read(member)
                        source_label = f"{uploaded.name} :: {member.filename}"
                        try:
                            raw, text, _ = parse_document(member.filename, member_data)
                            std, quality = standardize_transactions(raw, source_label)
                            standardized_frames.append(std)
                            all_text.append(text)
                            qualities.append(quality)
                            parsed_members += 1
                        except Exception as exc:
                            qualities.append({"source": source_label, "rows_raw": 0, "rows_valid": 0, "warnings": [str(exc)]})
                    if parsed_members == 0:
                        qualities.append({"source": uploaded.name, "rows_raw": 0, "rows_valid": 0, "warnings": ["Nessun file bancario supportato trovato nello ZIP"]})
            except Exception as exc:
                qualities.append({"source": uploaded.name, "rows_raw": 0, "rows_valid": 0, "warnings": [str(exc)]})
            continue
        try:
            raw, text, _ = parse_document(uploaded.name, data)
            std, quality = standardize_transactions(raw, uploaded.name)
            standardized_frames.append(std)
            all_text.append(text)
            qualities.append(quality)
        except Exception as exc:
            qualities.append({"source": uploaded.name, "rows_raw": 0, "rows_valid": 0, "warnings": [str(exc)]})

    if standardized_frames:
        merged = pd.concat(standardized_frames, ignore_index=True, sort=False)
        merged = merged.drop_duplicates(subset=["date", "description", "amount", "balance", "source"], keep="first")
        merged = merged.sort_values("date").reset_index(drop=True)
    else:
        merged = pd.DataFrame(columns=["date", "description", "amount", "balance", "source"])
    return merged, "\n".join(all_text), qualities, source_files


# -----------------------------------------------------------------------------
# CATEGORIZZAZIONE E KPI
# -----------------------------------------------------------------------------
def normalize_counterparty(description: str) -> str:
    text = clean_text(description).upper()
    text = re.sub(r"\b(?:TRN|CRO|ID|REF|RIF|SEPA|BONIFICO|ORDINANTE|BENEFICIARIO|PAGAMENTO|FATTURA|FT|N\.?|NR\.?|DEL)\b", " ", text)
    text = re.sub(r"\b[A-Z]{0,3}\d{5,}\b", " ", text)
    text = re.sub(r"\d{1,2}[/-]\d{1,2}[/-]\d{2,4}", " ", text)
    text = re.sub(r"[^A-ZÀ-ÖØ-Ý0-9 &'\.]+", " ", text)
    tokens = [t for t in text.split() if len(t) > 1]
    return " ".join(tokens[:7])[:80] or "CONTROPARTITA NON IDENTIFICATA"


def categorize_transaction(description: str, amount: float) -> str:
    text = clean_text(description).upper()
    for category, pattern in CATEGORY_PATTERNS:
        if re.search(pattern, text, flags=re.I):
            # L'anticipo fatture può essere entrata o uscita; mantiene categoria dedicata.
            return category
    if amount > 0:
        if re.search(r"\b(BONIFICO|ACCREDITO|INCASSO|PAGAMENTO FATTURA|POS|VERSAMENTO)\b", text):
            return "Incassi commerciali"
        if re.search(r"\b(FINANZIAMENTO|EROGAZIONE|PRESTITO|MUTUO)\b", text):
            return "Finanziamenti in entrata"
        return "Altre entrate"
    if re.search(r"\b(BONIFICO|PAGAMENTO|ADDEBITO)\b", text):
        return "Altre uscite operative"
    return "Altre uscite"


def month_span(start: pd.Timestamp, end: pd.Timestamp) -> int:
    return max(1, (end.year - start.year) * 12 + end.month - start.month + 1)


def recurring_counterparties(df: pd.DataFrame, sign: str, min_months: int = 3) -> tuple[set[str], float]:
    subset = df[df["amount"] > 0] if sign == "in" else df[df["amount"] < 0]
    if subset.empty:
        return set(), 0.0
    grouped = subset.groupby("counterparty").agg(
        total=("amount", lambda s: float(s.abs().sum())),
        months=("month", "nunique"),
        occurrences=("amount", "count"),
        mean_abs=("amount", lambda s: float(s.abs().mean())),
        std_abs=("amount", lambda s: float(s.abs().std(ddof=0))),
    )
    grouped["cv"] = grouped["std_abs"] / grouped["mean_abs"].replace(0, np.nan)
    recurring = grouped[(grouped["months"] >= min_months) & ((grouped["cv"] <= 0.55) | grouped["cv"].isna())]
    names = set(recurring.index)
    total_abs = float(subset["amount"].abs().sum())
    recurring_abs = float(subset[subset["counterparty"].isin(names)]["amount"].abs().sum())
    return names, recurring_abs / total_abs if total_abs else 0.0


def pv_annuity(payment: float, annual_rate: float = 0.08, months: int = 60) -> float:
    if payment <= 0:
        return 0.0
    r = annual_rate / 12
    if r == 0:
        return payment * months
    return payment * (1 - (1 + r) ** (-months)) / r


def clip(value: float, low: float, high: float) -> float:
    return max(low, min(high, value))


def score_band(score: float) -> tuple[str, str, float]:
    if score >= 85:
        return "A", "BASSO", 0.005
    if score >= 75:
        return "B", "BASSO", 0.012
    if score >= 65:
        return "C", "MEDIO-BASSO", 0.025
    if score >= 55:
        return "D", "MEDIO", 0.050
    if score >= 45:
        return "E", "MEDIO-ALTO", 0.090
    if score >= 35:
        return "F", "ALTO", 0.160
    return "G", "ALTO", 0.280


def semaforo_from_value(kind: str, value: float) -> str:
    if kind == "liquidity":
        return "VERDE" if value >= 0.08 else ("GIALLO" if value >= 0 else "ROSSO")
    if kind == "dscr":
        return "VERDE" if value >= 1.50 else ("GIALLO" if value >= 1.00 else "ROSSO")
    if kind == "redflags":
        return "VERDE" if value == 0 else ("GIALLO" if value <= 2 else "ROSSO")
    return "GIALLO"


def compute_analysis(
    transactions: pd.DataFrame,
    credit_limit: float = 0.0,
    requested_line: float = 0.0,
    manual_current_balance: Optional[float] = None,
) -> dict[str, Any]:
    if transactions.empty:
        raise ValueError("Nessun movimento valido da analizzare")

    df = transactions.copy()
    df["date"] = pd.to_datetime(df["date"])
    df = df.sort_values("date").reset_index(drop=True)
    df["description"] = df["description"].fillna("").map(clean_text)
    df["category"] = [categorize_transaction(d, a) for d, a in zip(df["description"], df["amount"])]
    df["counterparty"] = df["description"].map(normalize_counterparty)
    df["month"] = df["date"].dt.to_period("M").astype(str)

    start = df["date"].min()
    end = df["date"].max()
    months = month_span(start, end)

    inflows = df.loc[df["amount"] > 0, "amount"]
    outflows = -df.loc[df["amount"] < 0, "amount"]
    total_inflows = float(inflows.sum())
    total_outflows = float(outflows.sum())
    net_cashflow = total_inflows - total_outflows
    avg_monthly_inflows = total_inflows / months
    avg_monthly_outflows = total_outflows / months
    avg_monthly_net = net_cashflow / months

    monthly = df.groupby("month").agg(
        entrate=("amount", lambda s: float(s[s > 0].sum())),
        uscite=("amount", lambda s: float(-s[s < 0].sum())),
        netto=("amount", "sum"),
        movimenti=("amount", "count"),
    ).reset_index()

    # Saldi: usa saldo esplicito; in alternativa ricostruzione solo se saldo finale manuale.
    balances = df["balance"].dropna()
    balance_quality = "Saldo presente nell'estratto conto"
    if not balances.empty:
        current_balance = float(df.loc[df["balance"].notna(), "balance"].iloc[-1])
        avg_balance = float(balances.mean())
        min_balance = float(balances.min())
        max_balance = float(balances.max())
        df["effective_balance"] = df["balance"]
    elif manual_current_balance is not None:
        cumulative = df["amount"].cumsum()
        shift = float(manual_current_balance) - float(cumulative.iloc[-1])
        df["effective_balance"] = cumulative + shift
        current_balance = float(manual_current_balance)
        avg_balance = float(df["effective_balance"].mean())
        min_balance = float(df["effective_balance"].min())
        max_balance = float(df["effective_balance"].max())
        balance_quality = "Saldo ricostruito dal saldo finale indicato"
    else:
        df["effective_balance"] = np.nan
        current_balance = float("nan")
        avg_balance = float("nan")
        min_balance = float("nan")
        max_balance = float("nan")
        balance_quality = "Saldo non disponibile: KPI di giacenza parziali"

    if df["effective_balance"].notna().any():
        daily_last = df.dropna(subset=["effective_balance"]).groupby(df["date"].dt.date)["effective_balance"].last()
        negative_days = int((daily_last < 0).sum())
        overdraft_days = int((daily_last < -abs(credit_limit)).sum()) if credit_limit > 0 else negative_days
        volatility = float(daily_last.std(ddof=0)) if len(daily_last) > 1 else 0.0
        monthly_balance = df.dropna(subset=["effective_balance"]).groupby("month")["effective_balance"].mean().reset_index(name="saldo_medio")
    else:
        negative_days = 0
        overdraft_days = 0
        volatility = float("nan")
        monthly_balance = pd.DataFrame(columns=["month", "saldo_medio"])

    recurring_in_names, recurring_income_ratio = recurring_counterparties(df, "in", min_months=max(2, min(3, months)))
    recurring_out_names, fixed_cost_ratio = recurring_counterparties(df, "out", min_months=max(2, min(3, months)))

    pos = df[df["amount"] > 0]
    by_customer = pos.groupby("counterparty")["amount"].sum().sort_values(ascending=False)
    concentration_top3 = float(by_customer.head(3).sum() / total_inflows) if total_inflows else 0.0
    concentration_top5 = float(by_customer.head(5).sum() / total_inflows) if total_inflows else 0.0
    top_customers = [
        {"name": idx, "amount": float(value), "share": float(value / total_inflows) if total_inflows else 0.0}
        for idx, value in by_customer.head(5).items()
    ]

    category_out = df[df["amount"] < 0].groupby("category")["amount"].sum().abs().sort_values(ascending=False)
    category_in = df[df["amount"] > 0].groupby("category")["amount"].sum().sort_values(ascending=False)

    debt_service = float(df[(df["amount"] < 0) & (df["category"] == "Debito finanziario")]["amount"].abs().sum())
    payroll = float(df[(df["amount"] < 0) & (df["category"] == "Personale")]["amount"].abs().sum())
    taxes = float(df[(df["amount"] < 0) & (df["category"] == "Tasse e contributi")]["amount"].abs().sum())
    suppliers = float(df[(df["amount"] < 0) & (df["category"] == "Fornitori commerciali")]["amount"].abs().sum())
    utilities = float(df[(df["amount"] < 0) & (df["category"] == "Utenze e servizi")]["amount"].abs().sum())
    bank_costs = float(df[(df["amount"] < 0) & (df["category"] == "Spese bancarie")]["amount"].abs().sum())
    cash_withdrawals = float(df[(df["amount"] < 0) & (df["category"] == "Prelievi contante")]["amount"].abs().sum())
    related_parties = float(df[(df["amount"] < 0) & (df["category"] == "Giroconti/Parti correlate")]["amount"].abs().sum())
    invoice_advance_in = float(df[(df["amount"] > 0) & (df["category"] == "Anticipo fatture/Ri.Ba.")]["amount"].sum())
    invoice_advance_dependency = invoice_advance_in / total_inflows if total_inflows else 0.0

    non_debt_outflows = max(0.0, total_outflows - debt_service)
    cash_available_for_debt = total_inflows - non_debt_outflows
    dscr_cash = cash_available_for_debt / debt_service if debt_service > 0 else (3.0 if cash_available_for_debt > 0 else 0.0)
    dti = debt_service / total_inflows if total_inflows else 1.0
    opex_ratio = (suppliers + utilities + payroll + taxes) / total_inflows if total_inflows else 1.0
    payroll_ratio = payroll / total_inflows if total_inflows else 0.0
    payroll_monthly = (
        df[(df["amount"] < 0) & (df["category"] == "Personale")]
        .groupby("month")["amount"].sum().abs()
    )
    if payroll > 0 and not payroll_monthly.empty:
        payroll_coverage = min(1.0, payroll_monthly.size / max(months, 1))
        payroll_cv = float(payroll_monthly.std(ddof=0) / payroll_monthly.mean()) if payroll_monthly.mean() else 1.0
        payroll_stability = payroll_coverage * clip(1.0 - payroll_cv, 0.0, 1.0)
    else:
        payroll_stability = 0.0
    tax_pressure = taxes / total_inflows if total_inflows else 0.0

    avg_credit_utilization = float("nan")
    if credit_limit > 0 and df["effective_balance"].notna().any():
        used = (-df["effective_balance"]).clip(lower=0)
        avg_credit_utilization = float((used / credit_limit).mean())

    autonomy_months = 0.0
    if avg_monthly_outflows > 0 and not math.isnan(current_balance):
        autonomy_months = max(0.0, current_balance / avg_monthly_outflows)

    negative_months = monthly[monthly["netto"] < 0]
    burn_rate = float(-negative_months["netto"].mean()) if not negative_months.empty else 0.0

    # Trend ultimi tre mesi con regressione semplice.
    last3 = monthly.tail(3)
    if len(last3) >= 2:
        x = np.arange(len(last3))
        inflow_slope = float(np.polyfit(x, last3["entrate"], 1)[0])
        net_slope = float(np.polyfit(x, last3["netto"], 1)[0])
    else:
        inflow_slope = 0.0
        net_slope = 0.0

    # Red flags.
    red_flag_rows: list[dict[str, Any]] = []
    checks = [
        ("Insoluti/Ri.Ba./SDD respinti", df["category"].eq("Insoluti e respinti")),
        ("Pignoramenti/recupero crediti", df["category"].eq("Red flag - Recupero crediti/Pignoramenti")),
        ("Gioco e scommesse", df["category"].eq("Red flag - Gioco e scommesse")),
        ("Rateizzazioni/pressione riscossione", df["description"].str.upper().str.contains(r"ADER|RISCOSSIONE|RATEIZZAZIONE|ROTTAMAZIONE", regex=True)),
        ("Revoca/riduzione fidi", df["description"].str.upper().str.contains(r"REVOCA FID|REVOCA AFFIDAMENT|RIDUZIONE FID|RECESSO LINEA DI CREDITO", regex=True)),
        ("Transazioni con parti correlate", df["category"].eq("Giroconti/Parti correlate")),
        ("Paesi/giurisdizioni da verificare", df["description"].str.upper().str.contains("|".join(COUNTRY_RISK_KEYWORDS), regex=True)),
    ]
    for label, mask in checks:
        subset = df[mask]
        red_flag_rows.append({
            "label": label,
            "count": int(len(subset)),
            "amount": float(subset["amount"].abs().sum()),
            "status": "REGOLARE" if subset.empty else ("MONITORARE" if len(subset) <= 2 else "CRITICO"),
        })
    if credit_limit > 0:
        red_flag_rows.append({
            "label": "Sconfinamenti oltre fido",
            "count": overdraft_days,
            "amount": 0.0,
            "status": "REGOLARE" if overdraft_days == 0 else ("MONITORARE" if overdraft_days <= 5 else "CRITICO"),
        })
    else:
        red_flag_rows.append({
            "label": "Giorni con saldo negativo",
            "count": negative_days,
            "amount": 0.0,
            "status": "REGOLARE" if negative_days == 0 else ("MONITORARE" if negative_days <= 5 else "CRITICO"),
        })

    red_flags_count = sum(1 for row in red_flag_rows if row["status"] != "REGOLARE")
    critical_flags = sum(1 for row in red_flag_rows if row["status"] == "CRITICO")

    # Score interno 0-100.
    score_components: dict[str, float] = {}
    score_components["Cash flow"] = 20 * clip((avg_monthly_net / max(avg_monthly_inflows, 1) + 0.10) / 0.25, 0, 1)
    if math.isnan(avg_balance):
        score_components["Liquidita"] = 8.0
    else:
        liquidity_ratio = avg_balance / max(avg_monthly_outflows, 1)
        neg_penalty = clip(negative_days / max(1, (end - start).days + 1), 0, 1)
        score_components["Liquidita"] = 15 * clip(liquidity_ratio / 0.50 - neg_penalty, 0, 1)
    score_components["DSCR Cash"] = 20 * clip((dscr_cash - 0.75) / 1.25, 0, 1)
    score_components["DTI"] = 10 * clip((0.45 - dti) / 0.35, 0, 1)
    score_components["Concentrazione"] = 10 * clip((0.65 - concentration_top3) / 0.45, 0, 1)
    score_components["Ricorrenza"] = 10 * clip(recurring_income_ratio / 0.70, 0, 1)
    score_components["Red flags"] = max(0.0, 10.0 - red_flags_count * 2.5 - critical_flags * 2.0)
    if math.isnan(avg_credit_utilization):
        score_components["Uso fido"] = 3.0
    else:
        score_components["Uso fido"] = 5 * clip((0.95 - avg_credit_utilization) / 0.75, 0, 1)
    score = round(clip(sum(score_components.values()), 0, 100), 1)
    rating_class, risk_level, estimated_pd = score_band(score)

    # Capacita incrementale indicativa a 60 mesi, DSCR target 1,25.
    monthly_cash_available = cash_available_for_debt / months
    current_monthly_debt = debt_service / months
    additional_installment_capacity = max(0.0, monthly_cash_available / 1.25 - current_monthly_debt)
    suggested_line_cash = pv_annuity(additional_installment_capacity, annual_rate=0.08, months=60)
    turnover_cap = (total_inflows / months * 12) * 0.25
    suggested_line = round(max(0.0, min(suggested_line_cash, turnover_cap)), -2)

    if requested_line > 0:
        requested_installment = requested_line / max(pv_annuity(1.0, 0.08, 60), 1)
        post_dscr = monthly_cash_available / max(current_monthly_debt + requested_installment, 1)
        if post_dscr >= 1.50 and score >= 65 and red_flags_count <= 1:
            decision = "CONSIGLIATA"
        elif post_dscr >= 1.00 and score >= 50 and critical_flags == 0:
            decision = "VALUTABILE CON PRESIDI"
        else:
            decision = "NON CONSIGLIATA NELLA MISURA RICHIESTA"
    else:
        post_dscr = dscr_cash
        decision = "CAPACITA INDICATIVA CALCOLATA"

    liquidity_semaphore = semaforo_from_value("liquidity", avg_monthly_net / max(avg_monthly_inflows, 1))
    debt_semaphore = semaforo_from_value("dscr", dscr_cash)
    flags_semaphore = semaforo_from_value("redflags", red_flags_count)

    strengths = []
    weaknesses = []
    if avg_monthly_net > 0:
        strengths.append("cash flow medio mensile positivo")
    else:
        weaknesses.append("cash flow medio mensile negativo")
    if dscr_cash >= 1.5:
        strengths.append("DSCR Cash ampiamente superiore a 1")
    elif dscr_cash < 1:
        weaknesses.append("DSCR Cash inferiore a 1")
    if recurring_income_ratio >= 0.60:
        strengths.append("elevata ricorrenza degli incassi")
    elif recurring_income_ratio < 0.30:
        weaknesses.append("bassa ricorrenza degli incassi")
    if concentration_top3 > 0.50:
        weaknesses.append("concentrazione elevata sui primi tre clienti")
    else:
        strengths.append("concentrazione clienti contenuta")
    if red_flags_count:
        weaknesses.append(f"{red_flags_count} aree di allerta da approfondire")
    if invoice_advance_dependency > 0.35:
        weaknesses.append("dipendenza significativa da anticipo fatture/Ri.Ba.")

    summary = (
        f"L'analisi dei flussi dal {start.strftime('%d-%m-%Y')} al {end.strftime('%d-%m-%Y')} "
        f"attribuisce uno score interno di {score:.0f}/100 (classe {rating_class}, rischio {risk_level}). "
        f"Il cash flow netto del periodo è {euro(net_cashflow)} e il DSCR Cash è {num_it(dscr_cash)}. "
        + ("Punti di forza: " + ", ".join(strengths) + ". " if strengths else "")
        + ("Aree da monitorare: " + ", ".join(weaknesses) + "." if weaknesses else "Non emergono criticità rilevanti dai soli flussi disponibili.")
    )

    result = {
        "period_start": start.strftime("%Y-%m-%d"),
        "period_end": end.strftime("%Y-%m-%d"),
        "months": months,
        "transactions_count": int(len(df)),
        "total_inflows": total_inflows,
        "total_outflows": total_outflows,
        "net_cashflow": net_cashflow,
        "avg_monthly_inflows": avg_monthly_inflows,
        "avg_monthly_outflows": avg_monthly_outflows,
        "avg_monthly_net": avg_monthly_net,
        "current_balance": current_balance,
        "avg_balance": avg_balance,
        "min_balance": min_balance,
        "max_balance": max_balance,
        "negative_days": negative_days,
        "overdraft_days": overdraft_days,
        "balance_volatility": volatility,
        "balance_quality": balance_quality,
        "burn_rate": burn_rate,
        "autonomy_months": autonomy_months,
        "recurring_income_ratio": recurring_income_ratio,
        "fixed_cost_ratio": fixed_cost_ratio,
        "concentration_top3": concentration_top3,
        "concentration_top5": concentration_top5,
        "debt_service": debt_service,
        "dscr_cash": dscr_cash,
        "dti": dti,
        "opex_ratio": opex_ratio,
        "payroll": payroll,
        "payroll_ratio": payroll_ratio,
        "payroll_stability": payroll_stability,
        "taxes": taxes,
        "tax_pressure": tax_pressure,
        "suppliers": suppliers,
        "utilities": utilities,
        "bank_costs": bank_costs,
        "cash_withdrawals": cash_withdrawals,
        "related_parties": related_parties,
        "invoice_advance_dependency": invoice_advance_dependency,
        "credit_limit": credit_limit,
        "avg_credit_utilization": avg_credit_utilization,
        "requested_line": requested_line,
        "suggested_line": suggested_line,
        "post_dscr": post_dscr,
        "decision": decision,
        "score": score,
        "rating_class": rating_class,
        "risk_level": risk_level,
        "estimated_pd": estimated_pd,
        "score_components": score_components,
        "liquidity_semaphore": liquidity_semaphore,
        "debt_semaphore": debt_semaphore,
        "flags_semaphore": flags_semaphore,
        "red_flags_count": red_flags_count,
        "critical_flags": critical_flags,
        "red_flags": red_flag_rows,
        "top_customers": top_customers,
        "category_out": [{"category": idx, "amount": float(val)} for idx, val in category_out.items()],
        "category_in": [{"category": idx, "amount": float(val)} for idx, val in category_in.items()],
        "monthly": monthly.to_dict("records"),
        "monthly_balance": monthly_balance.to_dict("records"),
        "inflow_slope": inflow_slope,
        "net_slope": net_slope,
        "strengths": strengths,
        "weaknesses": weaknesses,
        "summary": summary,
        "methodology_note": (
            "Score, PD e importo suggerito sono stime interne basate esclusivamente sui dati caricati. "
            "Non sostituiscono istruttoria bancaria, Centrale Rischi, bilanci, visure, controlli AML/KYC o valutazioni MCC."
        ),
    }
    return result


def filter_period(df: pd.DataFrame, option: str, custom_start: Optional[date] = None, custom_end: Optional[date] = None) -> pd.DataFrame:
    if df.empty:
        return df
    work = df.copy()
    work["date"] = pd.to_datetime(work["date"])
    max_date = work["date"].max()
    if option.startswith("Ultimi "):
        months = int(re.search(r"\d+", option).group())
        min_date = max_date - pd.DateOffset(months=months) + pd.Timedelta(days=1)
        return work[(work["date"] >= min_date) & (work["date"] <= max_date)].copy()
    if option == "Periodo personalizzato" and custom_start and custom_end:
        return work[(work["date"].dt.date >= custom_start) & (work["date"].dt.date <= custom_end)].copy()
    return work


# -----------------------------------------------------------------------------
# GRAFICI PDF
# -----------------------------------------------------------------------------
def chart_monthly_flows(analysis: dict[str, Any]) -> io.BytesIO:
    monthly = pd.DataFrame(analysis["monthly"])
    fig, ax = plt.subplots(figsize=(8.2, 3.0))
    x = np.arange(len(monthly))
    width = 0.38
    ax.bar(x - width / 2, monthly["entrate"], width, label="Entrate", color="#16835A")
    ax.bar(x + width / 2, monthly["uscite"], width, label="Uscite", color="#B63131")
    ax.plot(x, monthly["netto"], marker="o", color="#0B1F3A", linewidth=1.8, label="Cash flow netto")
    ax.axhline(0, color="#7B8794", linewidth=0.8)
    ax.set_xticks(x)
    ax.set_xticklabels(monthly["month"], rotation=45, ha="right", fontsize=7)
    ax.tick_params(axis="y", labelsize=7)
    ax.grid(axis="y", alpha=0.2)
    handles, labels = ax.get_legend_handles_labels()
    balance_data = pd.DataFrame(analysis.get("monthly_balance", []))
    if not balance_data.empty:
        balance_map = dict(zip(balance_data["month"], balance_data["saldo_medio"]))
        balance_values = [balance_map.get(month, np.nan) for month in monthly["month"]]
        ax2 = ax.twinx()
        line2 = ax2.plot(x, balance_values, marker="s", linestyle="--", color="#B87333", linewidth=1.4, label="Saldo medio mensile")[0]
        ax2.tick_params(axis="y", labelsize=7, colors="#8A532A")
        handles.append(line2)
        labels.append("Saldo medio mensile")
    ax.legend(handles, labels, loc="upper left", ncol=2, fontsize=6.7, frameon=False)
    ax.set_title("Andamento mensile di entrate, uscite e cash flow", fontsize=10, fontweight="bold")
    fig.tight_layout()
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", dpi=170, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buffer.seek(0)
    return buffer


def chart_categories(analysis: dict[str, Any]) -> io.BytesIO:
    data = pd.DataFrame(analysis["category_out"]).head(7)
    fig, ax = plt.subplots(figsize=(7.8, 3.1))
    if data.empty:
        ax.text(0.5, 0.5, "Categorie di uscita non disponibili", ha="center", va="center")
        ax.axis("off")
    else:
        labels = [textwrap.shorten(x, width=28, placeholder="...") for x in data["category"]]
        y = np.arange(len(data))
        ax.barh(y, data["amount"], color="#12355B")
        ax.set_yticks(y)
        ax.set_yticklabels(labels, fontsize=7)
        ax.invert_yaxis()
        ax.tick_params(axis="x", labelsize=7)
        ax.grid(axis="x", alpha=0.2)
        ax.set_title("Principali categorie di uscita", fontsize=10, fontweight="bold")
    fig.tight_layout()
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", dpi=170, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buffer.seek(0)
    return buffer


# -----------------------------------------------------------------------------
# REPORT PDF
# -----------------------------------------------------------------------------
def hex_color(value: str) -> colors.Color:
    value = value.lstrip("#")
    return colors.Color(int(value[0:2], 16) / 255, int(value[2:4], 16) / 255, int(value[4:6], 16) / 255)


def pdf_status_color(status: str) -> colors.Color:
    return {
        "VERDE": hex_color(GREEN),
        "GIALLO": hex_color(AMBER),
        "ROSSO": hex_color(RED),
        "REGOLARE": hex_color(GREEN),
        "MONITORARE": hex_color(AMBER),
        "CRITICO": hex_color(RED),
    }.get(status, hex_color(MUTED))


def build_pdf_styles() -> dict[str, ParagraphStyle]:
    styles = getSampleStyleSheet()
    return {
        "title": ParagraphStyle(
            "FPTitle", parent=styles["Title"], fontName=PDF_FONT_BOLD,
            fontSize=18, leading=22, textColor=hex_color(NAVY), alignment=TA_LEFT, spaceAfter=6,
        ),
        "subtitle": ParagraphStyle(
            "FPSubtitle", parent=styles["Normal"], fontName=PDF_FONT,
            fontSize=9.5, leading=13, textColor=hex_color(MUTED), spaceAfter=8,
        ),
        "h1": ParagraphStyle(
            "FPH1", parent=styles["Heading1"], fontName=PDF_FONT_BOLD,
            fontSize=12.5, leading=15, textColor=hex_color(NAVY), spaceBefore=4, spaceAfter=7,
        ),
        "h2": ParagraphStyle(
            "FPH2", parent=styles["Heading2"], fontName=PDF_FONT_BOLD,
            fontSize=10, leading=12, textColor=hex_color(NAVY_2), spaceBefore=3, spaceAfter=5,
        ),
        "body": ParagraphStyle(
            "FPBody", parent=styles["BodyText"], fontName=PDF_FONT,
            fontSize=8.6, leading=12, textColor=hex_color(TEXT), spaceAfter=5,
        ),
        "small": ParagraphStyle(
            "FPSmall", parent=styles["BodyText"], fontName=PDF_FONT,
            fontSize=7.3, leading=9.5, textColor=hex_color(MUTED),
        ),
        "label": ParagraphStyle(
            "FPLabel", parent=styles["BodyText"], fontName=PDF_FONT_BOLD,
            fontSize=7.5, leading=9, textColor=hex_color(MUTED),
        ),
        "value": ParagraphStyle(
            "FPValue", parent=styles["BodyText"], fontName=PDF_FONT_BOLD,
            fontSize=8.4, leading=10, textColor=hex_color(TEXT),
        ),
        "center": ParagraphStyle(
            "FPCenter", parent=styles["BodyText"], fontName=PDF_FONT_BOLD,
            fontSize=8.2, leading=10, textColor=hex_color(TEXT), alignment=TA_CENTER,
        ),
        "right": ParagraphStyle(
            "FPRight", parent=styles["BodyText"], fontName=PDF_FONT,
            fontSize=8.2, leading=10, textColor=hex_color(TEXT), alignment=TA_RIGHT,
        ),
        "white": ParagraphStyle(
            "FPWhite", parent=styles["BodyText"], fontName=PDF_FONT_BOLD,
            fontSize=10, leading=12, textColor=colors.white, alignment=TA_CENTER,
        ),
    }


def para(text: Any, style: ParagraphStyle) -> Paragraph:
    value = clean_text(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return Paragraph(value, style)


def info_table(rows: list[tuple[str, str]], styles: dict[str, ParagraphStyle], widths: Optional[list[float]] = None) -> Table:
    data = [[para(k, styles["label"]), para(v, styles["value"])] for k, v in rows]
    table = Table(data, colWidths=widths or [4.1 * cm, 12.0 * cm], hAlign="LEFT")
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, -1), hex_color(LIGHT)),
        ("BOX", (0, 0), (-1, -1), 0.5, hex_color(MID)),
        ("INNERGRID", (0, 0), (-1, -1), 0.35, hex_color(MID)),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    return table


def metric_grid(metrics: list[tuple[str, str]], styles: dict[str, ParagraphStyle], columns: int = 4) -> Table:
    cells = []
    for label, value in metrics:
        cells.append([para(label, styles["label"]), para(value, styles["value"])])
    rows = []
    for i in range(0, len(cells), columns):
        group = cells[i:i + columns]
        while len(group) < columns:
            group.append([para("", styles["label"]), para("", styles["value"])])
        rows.append([Table([[cell[0]], [cell[1]]], colWidths=[3.7 * cm], style=TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), colors.white),
            ("BOX", (0, 0), (-1, -1), 0.5, hex_color(MID)),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("LEFTPADDING", (0, 0), (-1, -1), 5),
            ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ])) for cell in group])
    outer = Table(rows, colWidths=[4.0 * cm] * columns, hAlign="LEFT")
    outer.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 2),
        ("RIGHTPADDING", (0, 0), (-1, -1), 2),
        ("TOPPADDING", (0, 0), (-1, -1), 2),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
    ]))
    return outer


def page_header_footer(canvas: Any, doc: Any, logo_path: Optional[Path], footer_url: str, report_code: str) -> None:
    canvas.saveState()
    width, height = A4
    # Header line / brand.
    canvas.setFillColor(hex_color(NAVY))
    canvas.rect(0, height - 17 * mm, width, 17 * mm, stroke=0, fill=1)
    if logo_path and logo_path.exists():
        try:
            canvas.drawImage(str(logo_path), 15 * mm, height - 14.5 * mm, width=34 * mm, height=10.5 * mm, preserveAspectRatio=True, mask="auto")
        except Exception:
            pass
    canvas.setFont(PDF_FONT_BOLD, 10)
    canvas.setFillColor(colors.white)
    canvas.drawRightString(width - 15 * mm, height - 10.5 * mm, f"{BRAND} | {APP_TITLE}")

    # Footer.
    canvas.setStrokeColor(hex_color(MID))
    canvas.setLineWidth(0.5)
    canvas.line(15 * mm, 13 * mm, width - 15 * mm, 13 * mm)
    canvas.setFont(PDF_FONT, 7)
    canvas.setFillColor(hex_color(MUTED))
    canvas.drawString(15 * mm, 8.2 * mm, footer_url)
    canvas.drawCentredString(width / 2, 8.2 * mm, f"ID Report: {report_code}")
    canvas.drawRightString(width - 15 * mm, 8.2 * mm, f"Pagina {doc.page}")
    canvas.restoreState()


def generate_pdf(
    client: ClientInfo,
    analysis: dict[str, Any],
    bank_name: str,
    iban: str,
    period_option: str,
    report_code: str,
    consultant_name: str,
    footer_url: str,
    logo_path: Optional[Path],
) -> bytes:
    buffer = io.BytesIO()
    styles = build_pdf_styles()

    doc = BaseDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=15 * mm,
        rightMargin=15 * mm,
        topMargin=23 * mm,
        bottomMargin=18 * mm,
        title=f"{APP_TITLE} - {client.legal_name}",
        author=BRAND,
        subject="Valutazione creditizia basata su estratti conto",
    )
    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="normal")
    template = PageTemplate(
        id="FinancePlus",
        frames=[frame],
        onPage=lambda canv, d: page_header_footer(canv, d, logo_path, footer_url, report_code),
    )
    doc.addPageTemplates([template])
    story: list[Any] = []

    # PAGINA 1 - EXECUTIVE SUMMARY
    story.append(Paragraph("REPORT DI VALUTAZIONE CREDITIZIA PMI", styles["title"]))
    story.append(Paragraph("Analisi dei flussi di cassa e degli estratti conto bancari", styles["subtitle"]))
    story.append(info_table([
        ("Ragione sociale", client.legal_name or "Non disponibile"),
        ("Partita IVA / C.F.", " / ".join(filter(None, [client.vat_number, client.tax_code])) or "Non disponibile"),
        ("Sede", ", ".join(filter(None, [client.address, client.postal_code, client.city, client.province])) or "Non disponibile"),
        ("Settore ATECO", client.ateco or "Non disponibile"),
        ("Banca / IBAN", " - ".join(filter(None, [bank_name, iban])) or "Non disponibile"),
        ("Periodo analizzato", f"{datetime.fromisoformat(analysis['period_start']).strftime('%d-%m-%Y')} - {datetime.fromisoformat(analysis['period_end']).strftime('%d-%m-%Y')} ({analysis['months']} mesi; opzione: {period_option})"),
        ("Data elaborazione", now_display()),
        ("Consulente", consultant_name),
    ], styles))
    story.append(Spacer(1, 7))
    story.append(Paragraph("1. Executive summary e score synthesis", styles["h1"]))

    score_color = GREEN if analysis["score"] >= 75 else (AMBER if analysis["score"] >= 55 else RED)
    score_table = Table([
        [para(f"CLASSE {analysis['rating_class']}", styles["white"]), para(f"SCORE {analysis['score']:.0f}/100", styles["white"]), para(f"RISCHIO {analysis['risk_level']}", styles["white"])],
        [para("PD interna indicativa", styles["label"]), para(pct(analysis["estimated_pd"]), styles["value"]), para("Orizzonte 12 mesi - non regolamentare", styles["small"])],
    ], colWidths=[5.3 * cm, 5.3 * cm, 5.4 * cm])
    score_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), hex_color(score_color)),
        ("BACKGROUND", (0, 1), (-1, 1), hex_color(LIGHT)),
        ("BOX", (0, 0), (-1, -1), 0.7, hex_color(score_color)),
        ("INNERGRID", (0, 0), (-1, -1), 0.35, colors.white),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("TOPPADDING", (0, 0), (-1, -1), 7),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
    ]))
    story.append(score_table)
    story.append(Spacer(1, 7))

    sem_data = []
    for label, status in [
        ("Liquidita", analysis["liquidity_semaphore"]),
        ("Sostenibilita debito", analysis["debt_semaphore"]),
        ("Red flags", analysis["flags_semaphore"]),
    ]:
        sem_data.append([para(label, styles["label"]), para(status, styles["center"])])
    sem_table = Table([sem_data], colWidths=[5.35 * cm] * 3)
    for col_idx, (_, status) in enumerate([
        ("Liquidita", analysis["liquidity_semaphore"]),
        ("Sostenibilita debito", analysis["debt_semaphore"]),
        ("Red flags", analysis["flags_semaphore"]),
    ]):
        sem_table.setStyle(TableStyle([
            ("BOX", (col_idx, 0), (col_idx, 0), 1.0, pdf_status_color(status)),
            ("BACKGROUND", (col_idx, 0), (col_idx, 0), colors.white),
            ("TOPPADDING", (col_idx, 0), (col_idx, 0), 6),
            ("BOTTOMPADDING", (col_idx, 0), (col_idx, 0), 6),
        ]))
    story.append(sem_table)
    story.append(Spacer(1, 7))
    story.append(Paragraph("Commento automatico", styles["h2"]))
    story.append(Paragraph(analysis["summary"], styles["body"]))
    story.append(metric_grid([
        ("Entrate totali", euro(analysis["total_inflows"])),
        ("Uscite totali", euro(analysis["total_outflows"])),
        ("Cash flow netto", euro(analysis["net_cashflow"])),
        ("DSCR Cash", num_it(analysis["dscr_cash"])),
        ("Saldo attuale", euro(analysis["current_balance"])),
        ("Saldo medio", euro(analysis["avg_balance"])),
        ("Giorni in negativo", str(analysis["negative_days"])),
        ("Movimenti analizzati", str(analysis["transactions_count"])),
    ], styles, columns=4))

    # PAGINA 2 - FLUSSI E LIQUIDITA
    story.append(PageBreak())
    story.append(Paragraph("2. Dettaglio flussi di cassa e liquidita", styles["h1"]))
    flow_table_data = [[
        para("Metrica", styles["white"]), para("Media mensile", styles["white"]),
        para("Totale periodo", styles["white"]), para("Indicazione", styles["white"]),
    ]]
    trend_in = "In crescita" if analysis["inflow_slope"] > 0 else ("In calo" if analysis["inflow_slope"] < 0 else "Stabile")
    trend_net = "In miglioramento" if analysis["net_slope"] > 0 else ("In peggioramento" if analysis["net_slope"] < 0 else "Stabile")
    flow_table_data += [
        [para("Entrate", styles["body"]), para(euro(analysis["avg_monthly_inflows"]), styles["right"]), para(euro(analysis["total_inflows"]), styles["right"]), para(trend_in, styles["center"])],
        [para("Uscite", styles["body"]), para(euro(analysis["avg_monthly_outflows"]), styles["right"]), para(euro(analysis["total_outflows"]), styles["right"]), para("Gestione operativa", styles["center"])],
        [para("Cash flow netto", styles["body"]), para(euro(analysis["avg_monthly_net"]), styles["right"]), para(euro(analysis["net_cashflow"]), styles["right"]), para(trend_net, styles["center"])],
        [para("Servizio del debito", styles["body"]), para(euro(analysis["debt_service"] / analysis["months"]), styles["right"]), para(euro(analysis["debt_service"]), styles["right"]), para(f"DTI {pct(analysis['dti'])}", styles["center"])],
    ]
    flow_table = Table(flow_table_data, colWidths=[4.5 * cm, 3.7 * cm, 3.7 * cm, 4.2 * cm])
    flow_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), hex_color(NAVY)),
        ("BOX", (0, 0), (-1, -1), 0.5, hex_color(MID)),
        ("INNERGRID", (0, 0), (-1, -1), 0.35, hex_color(MID)),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, hex_color(LIGHT)]),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story.append(flow_table)
    story.append(Spacer(1, 6))
    story.append(Image(chart_monthly_flows(analysis), width=16.0 * cm, height=5.7 * cm))
    story.append(Spacer(1, 5))
    story.append(Paragraph("Analisi della liquidita e dei fidi", styles["h2"]))
    story.append(metric_grid([
        ("Saldo attuale aggregato", euro(analysis["current_balance"])),
        ("Giacenza media", euro(analysis["avg_balance"])),
        ("Saldo minimo", euro(analysis["min_balance"])),
        ("Saldo massimo", euro(analysis["max_balance"])),
        ("Volatilita saldo", euro(analysis["balance_volatility"])),
        ("Giorni saldo negativo", str(analysis["negative_days"])),
        ("Utilizzo medio fido", pct(analysis["avg_credit_utilization"])),
        ("Autonomia in stop incassi", f"{num_it(analysis['autonomy_months'], 1)} mesi"),
    ], styles, columns=4))
    story.append(Spacer(1, 5))
    story.append(Paragraph(analysis["balance_quality"], styles["small"]))

    # PAGINA 3 - KPI STRUTTURALI
    story.append(PageBreak())
    story.append(Paragraph("3. KPI strutturali, clienti e sostenibilita finanziaria", styles["h1"]))
    story.append(metric_grid([
        ("Ricorrenza incassi", pct(analysis["recurring_income_ratio"])),
        ("Concentrazione top 3", pct(analysis["concentration_top3"])),
        ("Concentrazione top 5", pct(analysis["concentration_top5"])),
        ("Costi fissi stimati", pct(analysis["fixed_cost_ratio"])),
        ("OpEx finanziario", pct(analysis["opex_ratio"])),
        ("Costo personale/entrate", pct(analysis["payroll_ratio"])),
        ("Stabilita payroll", pct(analysis["payroll_stability"])),
        ("Pressione fiscale", pct(analysis["tax_pressure"])),
        ("Dipendenza smobilizzo", pct(analysis["invoice_advance_dependency"])),
        ("DSCR Cash", num_it(analysis["dscr_cash"])),
        ("DTI", pct(analysis["dti"])),
        ("Burn rate", euro(analysis["burn_rate"])),
        ("Fido accordato", euro(analysis["credit_limit"])),
        ("Utilizzo medio fido", pct(analysis["avg_credit_utilization"])),
        ("Prelievi contante", euro(analysis["cash_withdrawals"])),
        ("Costi bancari", euro(analysis["bank_costs"])),
    ], styles, columns=4))
    story.append(Spacer(1, 6))
    story.append(Paragraph("Concentrazione degli incassi", styles["h2"]))
    top_data = [[para("Cliente/controparte", styles["white"]), para("Incassi", styles["white"]), para("Quota", styles["white"])]]
    for item in analysis["top_customers"]:
        top_data.append([
            para(item["name"], styles["body"]),
            para(euro(item["amount"]), styles["right"]),
            para(pct(item["share"]), styles["right"]),
        ])
    if len(top_data) == 1:
        top_data.append([para("Dati non disponibili", styles["body"]), para("-", styles["right"]), para("-", styles["right"])])
    top_table = Table(top_data, colWidths=[9.5 * cm, 3.7 * cm, 2.8 * cm])
    top_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), hex_color(NAVY)),
        ("BOX", (0, 0), (-1, -1), 0.5, hex_color(MID)),
        ("INNERGRID", (0, 0), (-1, -1), 0.35, hex_color(MID)),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, hex_color(LIGHT)]),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    story.append(top_table)
    story.append(Spacer(1, 6))
    story.append(Image(chart_categories(analysis), width=15.8 * cm, height=6.0 * cm))

    # PAGINA 4 - RED FLAGS E CONCLUSIONI
    story.append(PageBreak())
    story.append(Paragraph("4. Red flags, conclusioni e suggerimento di delibera", styles["h1"]))
    red_data = [[
        para("Verifica", styles["white"]), para("N.", styles["white"]),
        para("Importo", styles["white"]), para("Stato", styles["white"]),
    ]]
    for row in analysis["red_flags"]:
        red_data.append([
            para(row["label"], styles["body"]),
            para(str(row["count"]), styles["center"]),
            para(euro(row["amount"]), styles["right"]),
            para(row["status"], styles["center"]),
        ])
    red_table = Table(red_data, colWidths=[8.2 * cm, 1.5 * cm, 3.4 * cm, 3.0 * cm])
    red_style = [
        ("BACKGROUND", (0, 0), (-1, 0), hex_color(NAVY)),
        ("BOX", (0, 0), (-1, -1), 0.5, hex_color(MID)),
        ("INNERGRID", (0, 0), (-1, -1), 0.35, hex_color(MID)),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, hex_color(LIGHT)]),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]
    for idx, row in enumerate(analysis["red_flags"], start=1):
        red_style.append(("TEXTCOLOR", (3, idx), (3, idx), pdf_status_color(row["status"])))
    red_table.setStyle(TableStyle(red_style))
    story.append(red_table)
    story.append(Spacer(1, 8))
    story.append(Paragraph("Valutazione della nuova linea", styles["h2"]))
    decision_color = GREEN if analysis["decision"] == "CONSIGLIATA" else (AMBER if "VALUTABILE" in analysis["decision"] or "CALCOLATA" in analysis["decision"] else RED)
    decision_table = Table([
        [para(analysis["decision"], styles["white"])],
        [info_table([
            ("Linea richiesta", euro(analysis["requested_line"])),
            ("Linea indicativa suggerita", euro(analysis["suggested_line"])),
            ("DSCR post-operazione stimato", num_it(analysis["post_dscr"])),
            ("Presidio minimo", "DSCR target >= 1,25; verifica CR, bilanci, eventi e garanzie"),
        ], styles, widths=[5.2 * cm, 10.5 * cm])],
    ], colWidths=[16.1 * cm])
    decision_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), hex_color(decision_color)),
        ("BOX", (0, 0), (-1, -1), 0.8, hex_color(decision_color)),
        ("LEFTPADDING", (0, 1), (0, 1), 0),
        ("RIGHTPADDING", (0, 1), (0, 1), 0),
        ("TOPPADDING", (0, 0), (-1, 0), 7),
        ("BOTTOMPADDING", (0, 0), (-1, 0), 7),
    ]))
    story.append(decision_table)
    story.append(Spacer(1, 8))
    story.append(Paragraph("Conclusioni", styles["h2"]))
    story.append(Paragraph(analysis["summary"], styles["body"]))
    if analysis["strengths"]:
        story.append(Paragraph("Punti di forza: " + "; ".join(analysis["strengths"]) + ".", styles["body"]))
    if analysis["weaknesses"]:
        story.append(Paragraph("Azioni consigliate: verificare e presidiare " + "; ".join(analysis["weaknesses"]) + ".", styles["body"]))
    story.append(Spacer(1, 8))
    story.append(Paragraph("Avvertenza metodologica", styles["h2"]))
    story.append(Paragraph(analysis["methodology_note"], styles["small"]))
    story.append(Paragraph(
        "L'eventuale rilevazione di giurisdizioni o parole chiave sensibili è un filtro descrittivo e non equivale a verifica AML, sanzionatoria o reputazionale. Le causali devono essere validate documentalmente.",
        styles["small"],
    ))

    doc.build(story)
    return buffer.getvalue()


# -----------------------------------------------------------------------------
# ARCHIVIAZIONE REPORT
# -----------------------------------------------------------------------------
def generate_report_code(client_id: int, version: int) -> str:
    return f"CC-{datetime.now().strftime('%Y%m%d')}-{client_id:05d}-V{version}-{uuid.uuid4().hex[:5].upper()}"


def save_report_to_archive(
    client: ClientInfo,
    bank_name: str,
    iban: str,
    period_option: str,
    analysis: dict[str, Any],
    pdf_bytes: bytes,
    source_files: list[dict[str, Any]],
) -> tuple[int, str, Path]:
    client_id = upsert_client(client)
    period_start = analysis["period_start"]
    period_end = analysis["period_end"]
    version = count_matching_reports(client_id, bank_name, period_start, period_end) + 1
    report_code = generate_report_code(client_id, version)

    client_folder = CLIENTS_DIR / f"{client_id:05d}_{safe_filename(client.legal_name)}"
    report_folder = client_folder / "report_cc_flussi" / report_code
    source_folder = report_folder / "fonti"
    source_folder.mkdir(parents=True, exist_ok=True)

    for item in source_files:
        target = source_folder / safe_filename(Path(item["name"]).stem)
        target = target.with_suffix(Path(item["name"]).suffix.lower())
        if target.exists():
            target = target.with_name(f"{target.stem}_{item['hash'][:8]}{target.suffix}")
        target.write_bytes(item["bytes"])

    pdf_name = f"{safe_filename(client.legal_name)}_ANALISI_CC_FLUSSI_{datetime.now().strftime('%d-%m-%Y')}_{report_code}.pdf"
    pdf_path = report_folder / pdf_name
    pdf_path.write_bytes(pdf_bytes)

    record = {
        "client_id": client_id,
        "report_code": report_code,
        "version": version,
        "bank_name": bank_name,
        "iban": iban,
        "period_start": period_start,
        "period_end": period_end,
        "period_option": period_option,
        "source_names": json.dumps([x["name"] for x in source_files], ensure_ascii=False),
        "source_hash": combined_hash(source_files),
        "source_dir": str(source_folder),
        "pdf_path": str(pdf_path),
        "score": analysis["score"],
        "rating_class": analysis["rating_class"],
        "risk_level": analysis["risk_level"],
        "estimated_pd": analysis["estimated_pd"],
        "requested_line": analysis["requested_line"],
        "suggested_line": analysis["suggested_line"],
        "summary": analysis["summary"],
        "metrics_json": json.dumps(analysis, ensure_ascii=False, default=str),
        "created_at": now_iso(),
    }
    report_id = insert_report(record)
    return report_id, report_code, pdf_path


# -----------------------------------------------------------------------------
# STREAMLIT UI
# -----------------------------------------------------------------------------
def inject_css() -> None:
    st.markdown(
        f"""
        <style>
        :root {{ --navy:{NAVY}; --copper:{COPPER}; --light:{LIGHT}; }}
        .stApp {{ background: linear-gradient(180deg, #F7F9FC 0%, #EEF3F9 100%); }}
        [data-testid="stSidebar"] {{ background: linear-gradient(180deg, {NAVY} 0%, {NAVY_2} 100%); }}
        [data-testid="stSidebar"] * {{ color: white; }}
        .fp-hero {{
            background: linear-gradient(120deg, {NAVY} 0%, {NAVY_2} 72%, {COPPER} 160%);
            padding: 24px 28px; border-radius: 18px; color: white;
            box-shadow: 0 12px 30px rgba(11,31,58,.16); margin-bottom: 18px;
        }}
        .fp-hero h1 {{ margin:0; font-size:2rem; letter-spacing:-.02em; }}
        .fp-hero p {{ margin:.45rem 0 0; opacity:.88; }}
        .fp-card {{
            background:white; border:1px solid #DCE5F0; border-radius:14px;
            padding:16px 18px; box-shadow:0 6px 18px rgba(11,31,58,.06); margin-bottom:12px;
        }}
        .fp-kpi-label {{ color:{MUTED}; font-size:.8rem; font-weight:700; text-transform:uppercase; letter-spacing:.04em; }}
        .fp-kpi-value {{ color:{NAVY}; font-size:1.55rem; font-weight:800; margin-top:.2rem; }}
        .fp-badge {{ display:inline-block; border-radius:999px; padding:4px 10px; font-size:.78rem; font-weight:800; }}
        .fp-green {{ background:#E7F6EF; color:{GREEN}; }}
        .fp-amber {{ background:#FFF3D8; color:{AMBER}; }}
        .fp-red {{ background:#FDE8E8; color:{RED}; }}
        div.stButton > button {{ border-radius:10px; min-height:42px; font-weight:700; }}
        div.stDownloadButton > button {{ border-radius:10px; min-height:42px; font-weight:700; }}
        [data-testid="stMetricValue"] {{ color:{NAVY}; font-weight:800; }}
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_brand_header(page_title: str, subtitle: str) -> None:
    logo = find_logo_path()
    cols = st.columns([1, 5])
    with cols[0]:
        if logo:
            st.image(str(logo), use_container_width=True)
        else:
            st.markdown(f"<div class='fp-card' style='text-align:center;font-weight:900;color:{NAVY};'>FP<br>TECH</div>", unsafe_allow_html=True)
    with cols[1]:
        st.markdown(
            f"<div class='fp-hero'><h1>{page_title}</h1><p>{subtitle}</p></div>",
            unsafe_allow_html=True,
        )


def pdf_preview(pdf_bytes: bytes, height: int = 780) -> None:
    encoded = base64.b64encode(pdf_bytes).decode("ascii")
    html = f"""
    <iframe src="data:application/pdf;base64,{encoded}" width="100%" height="{height}px"
            style="border:1px solid #D9E2EF;border-radius:12px;background:white;"></iframe>
    """
    components.html(html, height=height + 20, scrolling=True)


def risk_badge(risk: str) -> str:
    css = "fp-green" if risk in {"BASSO", "MEDIO-BASSO"} else ("fp-amber" if risk in {"MEDIO", "MEDIO-ALTO"} else "fp-red")
    return f"<span class='fp-badge {css}'>{risk}</span>"


def page_dashboard() -> None:
    render_brand_header("Dashboard", "Controllo sintetico delle analisi di conto corrente e dei flussi aziendali")
    stats = dashboard_stats()
    cols = st.columns(4)
    values = [
        ("Clienti censiti", str(stats["clients"])),
        ("Report elaborati", str(stats["reports"])),
        ("Score medio", f"{stats['avg_score']:.1f}/100"),
        ("Report rischio alto", str(stats["high_risk"])),
    ]
    for col, (label, value) in zip(cols, values):
        with col:
            st.markdown(f"<div class='fp-card'><div class='fp-kpi-label'>{label}</div><div class='fp-kpi-value'>{value}</div></div>", unsafe_allow_html=True)

    st.subheader("Ultimi report elaborati")
    reports = list_reports_df().head(8)
    if reports.empty:
        st.info("Nessun report archiviato. Apri 'Inserisci E/C' per elaborare il primo estratto conto.")
    else:
        display = reports[["created_at", "legal_name", "bank_name", "period_start", "period_end", "score", "rating_class", "risk_level", "report_code"]].copy()
        display["created_at"] = pd.to_datetime(display["created_at"]).dt.strftime("%d-%m-%Y %H:%M")
        display["periodo"] = pd.to_datetime(display["period_start"]).dt.strftime("%d-%m-%Y") + " / " + pd.to_datetime(display["period_end"]).dt.strftime("%d-%m-%Y")
        display = display.rename(columns={
            "created_at": "Data", "legal_name": "Cliente", "bank_name": "Banca", "score": "Score",
            "rating_class": "Classe", "risk_level": "Rischio", "report_code": "ID Report",
        })
        st.dataframe(display[["Data", "Cliente", "Banca", "periodo", "Score", "Classe", "Rischio", "ID Report"]], use_container_width=True, hide_index=True)

    st.subheader("Flusso operativo")
    c1, c2 = st.columns(2)
    with c1:
        st.markdown("<div class='fp-card'><b>1. Inserisci E/C</b><br>Caricamento documenti, estrazione anagrafica, selezione periodo, KPI e anteprima PDF.</div>", unsafe_allow_html=True)
    with c2:
        st.markdown("<div class='fp-card'><b>2. Elenco Report E/C</b><br>Ricerca, anteprima, download e cancellazione controllata dei report archiviati.</div>", unsafe_allow_html=True)


def render_quality(qualities: list[dict[str, Any]]) -> None:
    with st.expander("Qualita dell'importazione e colonne riconosciute", expanded=False):
        for q in qualities:
            status = "OK" if q.get("rows_valid", 0) > 0 else "ATTENZIONE"
            st.markdown(f"**{q.get('source')}** - {status} - righe valide: {q.get('rows_valid', 0)}")
            cols = {k: v for k, v in q.items() if k.endswith("_col") and v}
            if cols:
                st.caption("Mappatura: " + ", ".join(f"{k.replace('_col','')}: {v}" for k, v in cols.items()))
            for warning in q.get("warnings", []):
                st.warning(warning)


def render_analysis_summary(analysis: dict[str, Any]) -> None:
    st.subheader("Risultato dell'analisi")
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Score interno", f"{analysis['score']:.0f}/100", f"Classe {analysis['rating_class']}")
    c2.metric("Cash flow netto", euro(analysis["net_cashflow"]))
    c3.metric("DSCR Cash", num_it(analysis["dscr_cash"]))
    c4.metric("Linea indicativa", euro(analysis["suggested_line"]))
    st.markdown(f"**Rischio:** {risk_badge(analysis['risk_level'])}", unsafe_allow_html=True)
    st.write(analysis["summary"])


def page_insert_statement(settings: dict[str, Any], selected_client: Optional[dict[str, Any]] = None) -> None:
    render_brand_header("Inserisci E/C", "Importa uno o piu estratti conto, analizza i flussi e genera il report PDF")

    uploaded_files = st.file_uploader(
        "Carica estratti conto o movimenti bancari",
        type=None,
        accept_multiple_files=True,
        help="Caricamento senza filtro di estensione: PDF, Word, Excel, CSV, Open Banking, immagini, ZIP e altri file. I formati non strutturati vengono comunque archiviati; l'analisi richiede movimenti leggibili.",
    )
    st.caption("Puoi selezionare file di qualsiasi estensione. PDF, Word, Excel, CSV, XML/OFX, immagini e ZIP vengono letti automaticamente quando possibile; i documenti non interpretabili restano comunque archiviabili come fonti.")

    if not uploaded_files:
        st.markdown("<div class='fp-card'><b>Nessun file caricato.</b><br>Seleziona almeno un estratto conto per avviare l'estrazione automatica.</div>", unsafe_allow_html=True)
        return

    with st.spinner("Lettura, normalizzazione e controllo dei movimenti..."):
        transactions, extracted_text, qualities, source_files = parse_uploaded_files(uploaded_files)
    render_quality(qualities)

    if transactions.empty:
        st.warning("I file sono stati caricati, ma non sono stati riconosciuti movimenti bancari utilizzabili. Puoi comunque correggere l'anagrafica; per elaborare il report aggiungi un estratto conto con data, causale e importo oppure un file Excel/CSV.")
        st.info("Nessun blocco sul formato: il documento e' accettato. L'elaborazione creditizia resta disabilitata finche' non sono disponibili movimenti validi.")
        return

    metadata = extract_metadata(extracted_text, [f.name for f in uploaded_files])
    if selected_client:
        metadata.legal_name = selected_client.get("legal_name") or metadata.legal_name
        metadata.vat_number = selected_client.get("vat_number") or metadata.vat_number
        metadata.address = selected_client.get("address") or metadata.address
        metadata.city = selected_client.get("city") or metadata.city
        metadata.province = selected_client.get("province") or metadata.province
        metadata.postal_code = selected_client.get("postal_code") or metadata.postal_code
    st.success(f"Movimenti riconosciuti: {len(transactions):,}".replace(",", "."))

    with st.expander("Anteprima movimenti normalizzati", expanded=False):
        preview = transactions.copy()
        preview["date"] = pd.to_datetime(preview["date"]).dt.strftime("%d-%m-%Y")
        preview["amount"] = preview["amount"].map(euro)
        preview["balance"] = preview["balance"].map(euro)
        st.dataframe(preview.head(300), use_container_width=True, hide_index=True)

    st.subheader("Anagrafica cliente")
    st.caption("Partita IVA e Codice fiscale sono facoltativi: il cliente e il report possono essere salvati anche senza questi dati.")
    c1, c2 = st.columns(2)
    legal_name = c1.text_input("Ragione sociale", value=metadata.legal_name, help="Facoltativa: in assenza di dati il sistema assegna un identificativo automatico.")
    vat_number = c2.text_input("Partita IVA", value=metadata.vat_number, max_chars=11)
    c3, c4 = st.columns(2)
    tax_code = c3.text_input("Codice fiscale", value=metadata.tax_code, max_chars=16)
    ateco = c4.text_input("Codice ATECO", value=metadata.ateco)
    address = st.text_input("Via / indirizzo", value=metadata.address)
    c5, c6, c7 = st.columns([2, 1, 1])
    city = c5.text_input("Comune", value=metadata.city)
    province = c6.text_input("Provincia", value=metadata.province, max_chars=2)
    postal_code = c7.text_input("CAP", value=metadata.postal_code, max_chars=5)

    st.subheader("Conto e periodo di analisi")
    c8, c9 = st.columns(2)
    bank_name = c8.text_input("Banca", value=metadata.bank_name)
    iban = c9.text_input("IBAN", value=metadata.iban)

    period_option = st.selectbox("Periodo da analizzare", ["Ultimi 6 mesi", "Ultimi 12 mesi", "Ultimi 24 mesi", "Tutto il periodo disponibile", "Periodo personalizzato"], index=1)
    custom_start = custom_end = None
    if period_option == "Periodo personalizzato":
        min_d = pd.to_datetime(transactions["date"]).min().date()
        max_d = pd.to_datetime(transactions["date"]).max().date()
        d1, d2 = st.columns(2)
        custom_start = d1.date_input("Data iniziale", value=min_d, min_value=min_d, max_value=max_d)
        custom_end = d2.date_input("Data finale", value=max_d, min_value=min_d, max_value=max_d)
        if custom_start > custom_end:
            st.error("La data iniziale non puo essere successiva alla data finale.")
            return

    c10, c11, c12 = st.columns(3)
    credit_limit = c10.number_input("Fido di conto accordato (€)", min_value=0.0, value=float(settings.get("default_credit_limit", 0.0)), step=1000.0)
    requested_line = c11.number_input("Nuova linea richiesta (€)", min_value=0.0, value=float(settings.get("default_requested_line", 0.0)), step=5000.0)
    use_manual_balance = c12.checkbox("Indica saldo finale", value=False)
    manual_current_balance: Optional[float] = None
    if use_manual_balance:
        manual_current_balance = st.number_input("Saldo finale/attuale (€)", value=0.0, step=1000.0)

    filtered = filter_period(transactions, period_option, custom_start, custom_end)
    if filtered.empty:
        st.error("Il periodo selezionato non contiene movimenti.")
        return
    st.caption(f"Periodo effettivo: {pd.to_datetime(filtered['date']).min().strftime('%d-%m-%Y')} - {pd.to_datetime(filtered['date']).max().strftime('%d-%m-%Y')} | {len(filtered)} movimenti")

    if st.button("Elabora report", type="primary", use_container_width=True):
        with st.spinner("Calcolo KPI, score, red flags e composizione del PDF..."):
            fallback_name = f"CLIENTE_NON_IDENTIFICATO_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            client = ClientInfo(
                legal_name=legal_name.strip() or fallback_name, vat_number=re.sub(r"\D", "", vat_number),
                tax_code=tax_code.strip().upper(), address=address.strip(), city=city.strip(),
                province=province.strip().upper(), postal_code=postal_code.strip(), ateco=ateco.strip(),
                bank_name=bank_name.strip(), iban=re.sub(r"\s+", "", iban).upper(),
            )
            analysis = compute_analysis(filtered, credit_limit, requested_line, manual_current_balance)
            temp_code = f"ANTEPRIMA-{datetime.now().strftime('%Y%m%d-%H%M%S')}"
            pdf_bytes = generate_pdf(
                client=client, analysis=analysis, bank_name=bank_name.strip(), iban=client.iban,
                period_option=period_option, report_code=temp_code,
                consultant_name=settings.get("consultant_name", BRAND),
                footer_url=settings.get("footer_url", FOOTER_URL_DEFAULT),
                logo_path=find_logo_path(),
            )
            st.session_state["cc_report_payload"] = {
                "client": asdict(client),
                "analysis": analysis,
                "bank_name": bank_name.strip(),
                "iban": client.iban,
                "period_option": period_option,
                "pdf_bytes": pdf_bytes,
                "source_files": source_files,
                "transactions": filtered,
                }
            st.success("Report elaborato. Ora puoi visualizzare l'anteprima o salvarlo nell'archivio.")

    payload = st.session_state.get("cc_report_payload")
    if payload:
        analysis = payload["analysis"]
        render_analysis_summary(analysis)
        b1, b2 = st.columns(2)
        with b1:
            show_preview = st.button("Anteprima report PDF", use_container_width=True)
        with b2:
            save_click = st.button("Salva report nell'elenco", type="primary", use_container_width=True)

        if show_preview:
            pdf_preview(payload["pdf_bytes"])
        st.download_button(
            "Scarica anteprima PDF",
            data=payload["pdf_bytes"],
            file_name=f"ANTEPRIMA_{safe_filename(payload['client']['legal_name'])}_ANALISI_CC_FLUSSI.pdf",
            mime="application/pdf",
            use_container_width=True,
        )

        if save_click:
            client = ClientInfo(**payload["client"])
            # Rigenera con ID definitivo e versione corretta.
            client_id = upsert_client(client)
            version = count_matching_reports(client_id, payload["bank_name"], analysis["period_start"], analysis["period_end"]) + 1
            definitive_code = generate_report_code(client_id, version)
            definitive_pdf = generate_pdf(
                client=client, analysis=analysis, bank_name=payload["bank_name"], iban=payload["iban"],
                period_option=payload["period_option"], report_code=definitive_code,
                consultant_name=settings.get("consultant_name", BRAND), footer_url=settings.get("footer_url", FOOTER_URL_DEFAULT),
                logo_path=find_logo_path(),
            )
            # Salvataggio manuale per mantenere il codice definitivo appena creato.
            client_folder = CLIENTS_DIR / f"{client_id:05d}_{safe_filename(client.legal_name)}"
            report_folder = client_folder / "report_cc_flussi" / definitive_code
            source_folder = report_folder / "fonti"
            source_folder.mkdir(parents=True, exist_ok=True)
            for item in payload["source_files"]:
                target = source_folder / (safe_filename(Path(item["name"]).stem) + Path(item["name"]).suffix.lower())
                if target.exists():
                    target = target.with_name(f"{target.stem}_{item['hash'][:8]}{target.suffix}")
                target.write_bytes(item["bytes"])
            pdf_name = f"{safe_filename(client.legal_name)}_ANALISI_CC_FLUSSI_{datetime.now().strftime('%d-%m-%Y')}_{definitive_code}.pdf"
            pdf_path = report_folder / pdf_name
            pdf_path.write_bytes(definitive_pdf)
            record = {
                "client_id": client_id,
                "report_code": definitive_code,
                "version": version,
                "bank_name": payload["bank_name"],
                "iban": payload["iban"],
                "period_start": analysis["period_start"],
                "period_end": analysis["period_end"],
                "period_option": payload["period_option"],
                "source_names": json.dumps([x["name"] for x in payload["source_files"]], ensure_ascii=False),
                "source_hash": combined_hash(payload["source_files"]),
                "source_dir": str(source_folder),
                "pdf_path": str(pdf_path),
                "score": analysis["score"],
                "rating_class": analysis["rating_class"],
                "risk_level": analysis["risk_level"],
                "estimated_pd": analysis["estimated_pd"],
                "requested_line": analysis["requested_line"],
                "suggested_line": analysis["suggested_line"],
                "summary": analysis["summary"],
                "metrics_json": json.dumps(analysis, ensure_ascii=False, default=str),
                "created_at": now_iso(),
            }
            insert_report(record)
            st.session_state["cc_report_payload"]["pdf_bytes"] = definitive_pdf
            st.success(f"Report salvato: {definitive_code} - versione {version}. L'anagrafica cliente è stata aggiornata senza creare duplicati.")


def page_reports() -> None:
    render_brand_header("Elenco Report E/C", "Consulta, filtra, visualizza e gestisci i report di analisi dei flussi")
    reports = list_reports_df()
    if reports.empty:
        st.info("Nessun report archiviato.")
        return

    f1, f2, f3 = st.columns(3)
    client_options = ["Tutti"] + sorted(reports["legal_name"].dropna().unique().tolist())
    bank_options = ["Tutte"] + sorted([x for x in reports["bank_name"].dropna().unique().tolist() if x])
    risk_options = ["Tutti"] + sorted(reports["risk_level"].dropna().unique().tolist())
    client_filter = f1.selectbox("Cliente", client_options)
    bank_filter = f2.selectbox("Banca", bank_options)
    risk_filter = f3.selectbox("Rischio", risk_options)

    filtered = reports.copy()
    if client_filter != "Tutti":
        filtered = filtered[filtered["legal_name"] == client_filter]
    if bank_filter != "Tutte":
        filtered = filtered[filtered["bank_name"] == bank_filter]
    if risk_filter != "Tutti":
        filtered = filtered[filtered["risk_level"] == risk_filter]

    search = st.text_input("Ricerca libera", placeholder="Ragione sociale, Partita IVA, banca o ID report")
    if search:
        needle = search.lower()
        mask = filtered[["legal_name", "vat_number", "bank_name", "report_code", "summary"]].fillna("").astype(str).apply(lambda col: col.str.lower().str.contains(needle, regex=False)).any(axis=1)
        filtered = filtered[mask]

    table = filtered.copy()
    table["Data elaborazione"] = pd.to_datetime(table["created_at"]).dt.strftime("%d-%m-%Y %H:%M")
    table["Periodo"] = pd.to_datetime(table["period_start"]).dt.strftime("%d-%m-%Y") + " / " + pd.to_datetime(table["period_end"]).dt.strftime("%d-%m-%Y")
    table["Punteggio"] = table["score"].round(1).astype(str) + "/100 - " + table["rating_class"].fillna("")
    table["Descrizione breve"] = table["summary"].fillna("").map(lambda s: textwrap.shorten(s, width=140, placeholder="..."))
    display_cols = ["Data elaborazione", "legal_name", "bank_name", "Periodo", "Punteggio", "risk_level", "version", "report_code", "Descrizione breve"]
    st.dataframe(
        table[display_cols].rename(columns={
            "legal_name": "Cliente", "bank_name": "Banca", "risk_level": "Rischio",
            "version": "Versione", "report_code": "ID Report",
        }),
        use_container_width=True,
        hide_index=True,
    )

    if filtered.empty:
        st.warning("Nessun report corrisponde ai filtri.")
        return

    labels = {
        f"{row['created_at']} | {row['legal_name']} | {row['bank_name'] or 'Banca n.d.'} | {row['score']:.0f}/100 | {row['report_code']}": int(row["id"])
        for _, row in filtered.iterrows()
    }
    selection = st.selectbox("Seleziona il report da aprire", list(labels.keys()))
    selected_id = labels[selection]
    row = filtered[filtered["id"] == selected_id].iloc[0]

    st.markdown(
        f"<div class='fp-card'><b>{row['legal_name']}</b><br>"
        f"ID {row['report_code']} - Score {row['score']:.0f}/100 - Classe {row['rating_class']} - "
        f"{risk_badge(row['risk_level'])}<br><span style='color:{MUTED}'>{row['summary']}</span></div>",
        unsafe_allow_html=True,
    )

    path = Path(row["pdf_path"])
    if not path.exists():
        st.error("Il PDF archiviato non è più presente sul disco.")
        return
    pdf_bytes = path.read_bytes()
    b1, b2 = st.columns(2)
    show = b1.button("Genera anteprima", use_container_width=True)
    b2.download_button("Scarica report PDF", data=pdf_bytes, file_name=path.name, mime="application/pdf", use_container_width=True)
    if show:
        pdf_preview(pdf_bytes)

    with st.expander("Elimina report", expanded=False):
        confirm = st.checkbox(f"Confermo l'eliminazione definitiva del report {row['report_code']}", key=f"confirm_{selected_id}")
        if st.button("Elimina report selezionato", type="secondary", disabled=not confirm, key=f"delete_{selected_id}"):
            delete_report(selected_id)
            st.success("Report eliminato dall'elenco e dall'archivio documentale.")
            st.rerun()


def page_settings(settings: dict[str, Any]) -> dict[str, Any]:
    render_brand_header("Impostazioni", "Brand, piè di pagina e valori predefiniti dell'applicazione")
    st.subheader("Logo FinancePlus.Tech")
    logo = find_logo_path()
    if logo:
        st.image(str(logo), width=260)
        st.caption(f"Logo attivo: {logo.name}")
    uploaded_logo = st.file_uploader("Carica o sostituisci il logo", type=["png", "jpg", "jpeg"])
    if uploaded_logo and st.button("Salva logo"):
        data = uploaded_logo.getvalue()
        if PILImage is not None:
            image = PILImage.open(io.BytesIO(data)).convert("RGBA")
            image.thumbnail((1800, 900))
            image.save(ASSETS_DIR / "logo_financeplus.png", format="PNG", optimize=True)
        else:
            (ASSETS_DIR / "logo_financeplus.png").write_bytes(data)
        st.success("Logo salvato. Sarà usato nell'app e nei report PDF.")
        st.rerun()

    st.subheader("Dati report")
    footer_url = st.text_input("Piè di pagina", value=settings.get("footer_url", FOOTER_URL_DEFAULT))
    consultant_name = st.text_input("Consulente / struttura", value=settings.get("consultant_name", BRAND))
    c1, c2 = st.columns(2)
    default_credit_limit = c1.number_input("Fido predefinito (€)", min_value=0.0, value=float(settings.get("default_credit_limit", 0.0)), step=1000.0)
    default_requested_line = c2.number_input("Linea richiesta predefinita (€)", min_value=0.0, value=float(settings.get("default_requested_line", 0.0)), step=5000.0)
    if st.button("Salva impostazioni", type="primary"):
        settings = {
            "footer_url": footer_url.strip() or FOOTER_URL_DEFAULT,
            "consultant_name": consultant_name.strip() or BRAND,
            "default_credit_limit": default_credit_limit,
            "default_requested_line": default_requested_line,
        }
        save_settings(settings)
        st.success("Impostazioni salvate.")

    st.subheader("Archivio locale")
    st.code(str(DATA_DIR))
    st.caption("Database SQLite e cartelle cliente vengono creati automaticamente accanto al file .py. Su Streamlit Cloud utilizzare uno storage persistente esterno per evitare la perdita dei dati al riavvio del container.")
    return settings


def main(embedded: bool = False, selected_client: Optional[dict[str, Any]] = None) -> None:
    if not embedded:
        st.set_page_config(page_title=f"{APP_TITLE} | {BRAND}", page_icon="📊", layout="wide", initial_sidebar_state="expanded")
    ensure_directories()
    init_db()
    inject_css()
    settings = load_settings()

    if embedded:
        st.markdown("---")
        st.markdown("## 📊 Analisi CC / Flussi")
        if selected_client:
            st.caption(
                f"Cliente collegato: {selected_client.get('legal_name', '')}"
                + (f" · P.IVA {selected_client.get('vat_number')}" if selected_client.get('vat_number') else "")
            )
        page = st.radio(
            "Navigazione Analisi CC",
            ["Dashboard", "Inserisci E/C", "Elenco Report E/C", "Impostazioni"],
            horizontal=True,
            key="cc_flussi_embedded_page",
        )
    else:
        with st.sidebar:
            logo = find_logo_path()
            if logo:
                st.image(str(logo), use_container_width=True)
            st.markdown(f"### {APP_TITLE}")
            st.caption(f"{BRAND} - versione {APP_VERSION}")
            page = st.radio("Navigazione", ["Dashboard", "Inserisci E/C", "Elenco Report E/C", "Impostazioni"], label_visibility="collapsed")
            st.divider()
            st.caption("Analisi gestionale dei flussi. Validare sempre dati, causali e documentazione prima di assumere decisioni di credito.")

    if page == "Dashboard":
        page_dashboard()
    elif page == "Inserisci E/C":
        page_insert_statement(settings, selected_client=selected_client)
    elif page == "Elenco Report E/C":
        page_reports()
    else:
        page_settings(settings)


if __name__ == "__main__":
    main()
